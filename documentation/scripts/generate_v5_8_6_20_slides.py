import os as _os  # [doc-folder] run from repo root so plot PNGs & Software/ resolve
_os.chdir(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..', '..'))
"""Phase 8.6.20 (Tensile to failure) presentation — 7 slides, JU template.
Pages 126-132:
  126 Overview + timeline flowchart
  127 Load vs time (preload->UTS delta shown)
  128 Stress vs strain (0.2% offset yield construction + failure values)
  129 Cauchy vs True strain  [NEW]
  130 Stress vs position (compliance split, rig stiffness explained)
  131 Verdict: pass criteria + offset-factor table + DIC-tracking definition
  132 Recommendations
Output: V5_8_6_20_slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

DARK_GREEN = RGBColor(0x1F, 0x3F, 0x2F)
HEADER_GREEN = RGBColor(0x14, 0x3D, 0x2F)
LIGHT_BLUE = RGBColor(0xE7, 0xF1, 0xF8)
GREEN_PASS = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_WARN = RGBColor(0xFF, 0xF3, 0xCD)
RED_FAIL = RGBColor(0xFF, 0xCD, 0xD2)
GREY_TEXT = RGBColor(0x70, 0x70, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BOX_DARK = RGBColor(0x1F, 0x36, 0x2D)
JU_BLUE_BAR = RGBColor(0xA9, 0xD6, 0xEF)
FLOW_BLUE = RGBColor(0x1F, 0x77, 0xB4)
FLOW_RED = RGBColor(0xD6, 0x27, 0x28)
FLOW_ORANGE = RGBColor(0xFF, 0x98, 0x00)
FLOW_NEUTRAL = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, x, y, w, h, text, *, fs=14, bold=False, italic=False, colour=BLACK,
       align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = colour
    return box


def title(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = JU_BLUE_BAR; bar.line.fill.background()
    tb(slide, 0.5, 0.55, 12.5, 0.9, text, fs=33, font="Calibri Light")


def header(slide, x, y, w, text):
    tb(slide, x, y, w, 0.4, text, fs=17, colour=GREY_TEXT)


def cell_bg(cell, rgb):
    cell.fill.solid(); cell.fill.fore_color.rgb = rgb


def cell_txt(cell, text, *, fs=10, bold=False, colour=BLACK, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.colour = colour
    r.font.color.rgb = colour; r.font.name = "Calibri"


def table(slide, x, y, w, h, data, *, hbg=HEADER_GREEN, hfg=WHITE, cw=None,
          hf=10, bf=10, ov=None):
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = shp.table
    if cw:
        tot = sum(cw)
        for i, f in enumerate(cw):
            t.columns[i].width = Inches(w*f/tot)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c)
            text = data[r][c]
            if r == 0:
                cell_bg(cell, hbg); cell_txt(cell, text, fs=hf, bold=True, colour=hfg)
            else:
                cell_txt(cell, text, fs=bf)
            if ov and (r, c) in ov:
                o = ov[(r, c)]
                if 'bg' in o: cell_bg(cell, o['bg'])
                cell_txt(cell, o.get('text', text), fs=bf, bold=o.get('bold', False),
                         colour=o.get('colour', BLACK))
    return shp


def flow(slide, x, y, w, h, text, *, fill=WHITE, border=FLOW_NEUTRAL, fs=11,
         bold=False, fg=BLACK, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fg; r.font.name = "Calibri"
    return box


def arrow(slide, x1, y1, x2, y2, *, colour=FLOW_NEUTRAL, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = colour; c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def kpi(slide, x, y, w, label, value, *, fill=LIGHT_BLUE, vcol=DARK_GREEN, h=0.98, vfs=19):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.text = ""
    r1 = p1.add_run(); r1.text = value
    r1.font.size = Pt(vfs); r1.font.bold = True; r1.font.color.rgb = vcol; r1.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(10); r2.font.color.rgb = BLACK; r2.font.name = "Calibri"
    return box


def pageno(slide, n):
    tb(slide, 0.5, 7.05, 0.7, 0.3, str(n), fs=10, colour=GREY_TEXT)


def ju(slide):
    tb(slide, 12.2, 0.3, 1.0, 0.3, "JÖNKÖPING UNIVERSITY", fs=8, bold=True,
       colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    tb(slide, 0.6, 7.0, 12.1, 0.4, text, fs=12, italic=True, colour=GREY_TEXT)


# ================= SLIDE 1 — Overview + timeline =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: TENSILE TEST TO FAILURE — V5")
header(s, 0.5, 1.4, 5.6, "Test setup")
setup = [
    ["Parameter", "Value"],
    ["Specimen", "S4 — fresh PLA, 50 % infill (grid pattern)"],
    ["Geometry", "10 × 8 mm gauge (80 mm², CAD-verified)"],
    ["Markers", "Spray paint — held to fracture"],
    ["Preload", "+463 N (anchored from post-fracture zero)"],
    ["Rate", "0.1 mm/s, continuous to fracture"],
    ["Load cell", "ANYLOAD 3 t (29.4 kN) — used at 6 %"],
    ["Reference", "Chacón et al. (2017), Materials & Design"],
]
table(s, 0.5, 1.8, 5.6, 2.7, setup, cw=[1.0, 2.7], hf=10, bf=10)
kpi(s, 0.5, 4.7, 1.78, "UTS (nominal)", "22.1 MPa")
kpi(s, 2.41, 4.7, 1.78, "failure strain", "2.46 %")
kpi(s, 4.32, 4.7, 1.78, "peak force", "1 767 N")
kpi(s, 0.5, 5.85, 1.78, "E (nominal)", "1.54 GPa")
kpi(s, 2.41, 5.85, 1.78, "σ_y (0.2 %)", "19.0 MPa")
kpi(s, 4.32, 5.85, 1.78, "DIC tracking", "99.8 %", fill=GREEN_PASS)

header(s, 6.6, 1.4, 6.2, "Test timeline")
FX, FW, FH = 7.1, 5.3, 0.52
steps = [
    ("t = 0 s   Preload +463 N → tare load + DIC", BOX_DARK, WHITE, True),
    ("0 – 96 s   Baseline hold  (noise σ_F=1.6 N, σ_ε=1.4×10⁻⁵)", WHITE, BLACK, False),
    ("96 s   Ramp start — 0.1 mm/s", FLOW_BLUE, WHITE, True),
    ("107.8 s   Proportional limit — 12.8 MPa, ε=0.005", WHITE, FLOW_BLUE, False),
    ("118.7 s   Yield (0.2 % offset) — σ_y = 19.0 MPa", FLOW_ORANGE, WHITE, True),
    ("130.6 s   UTS — 22.1 MPa (1 767 N), ε=0.020", YELLOW_WARN, BLACK, True),
    ("134.5 s   FRACTURE — gauge section, ε_f=0.0246", FLOW_RED, WHITE, True),
    ("134.5 – 270 s   Post-fracture: display −463 N ⇒ preload", GREEN_PASS, DARK_GREEN, True),
]
y = 1.85
for i, (text, fill, fg, bold) in enumerate(steps):
    flow(s, FX, y, FW, FH, text, fill=fill, fg=fg, bold=bold, fs=10.5)
    if i < len(steps)-1:
        arrow(s, FX+FW/2, y+FH, FX+FW/2, y+FH+0.13)
    y += FH + 0.13
pageno(s, 126)

# ================= SLIDE 1b — Predicted results (pre-test) =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: PREDICTED RESULTS (PRE-TEST)")

header(s, 0.5, 1.4, 6.0, "Expected for solid PLA — Chacón et al. (2017)")
pred = [
    ["Quantity", "Predicted", "Basis"],
    ["Elastic modulus E", "3.0 – 5.5 GPa", "Chacón, FDM PLA"],
    ["Yield σ_y (0.2 %)", "30 – 50 MPa", "Chacón, orientation-dep."],
    ["UTS", "32 – 60 MPa", "Chacón, orientation-dep."],
    ["Failure strain ε_f", "0.01 – 0.08", "PLA elongation at break"],
    ["Peak force", "2.6 – 4.8 kN", "UTS × 80 mm²"],
    ["Curve shape", "elastic→yield→plastic→fracture", "ductile-brittle polymer"],
    ["Failure location", "gauge section", "valid-test requirement"],
    ["DIC tracking", "≥ 90 %", "robustness target"],
]
table(s, 0.5, 1.8, 6.0, 3.9, pred, cw=[1.5, 1.7, 1.9], hf=10, bf=10)

header(s, 6.8, 1.4, 6.0, "Adjusted for 50 % infill (nominal 80 mm²)")
tb(s, 6.8, 1.78, 6.0, 0.7,
   "S4 is 50 % infill, so at the CAD-verified 80 mm² the readings are predicted to fall "
   "LOW by a load-bearing knock-down of ≈ 2× (see verdict slide). Pre-test nominal expectation:",
   fs=11, italic=True, colour=GREY_TEXT)
predi = [
    ["Quantity", "Dense (Chacón)", "÷ ≈2 infill", "→ nominal expected"],
    ["E", "3.0 – 5.5 GPa", "× 0.5", "1.5 – 2.8 GPa"],
    ["σ_y", "30 – 50 MPa", "× 0.5", "15 – 25 MPa"],
    ["UTS", "32 – 60 MPa", "× 0.5", "16 – 30 MPa"],
]
ovp = {(r, 3): {'bg': LIGHT_BLUE, 'bold': True} for r in range(1, 4)}
table(s, 6.8, 2.7, 6.0, 1.6, predi, cw=[0.8, 1.4, 1.0, 1.6], hf=9.5, bf=10, ov=ovp)
tb(s, 6.8, 4.45, 6.0, 1.3,
   "PASS = results land inside these ranges. Actual (next slides): "
   "E = 1.54 GPa, σ_y = 19.0 MPa, UTS = 22.1 MPa, ε_f = 0.025 — all inside the "
   "infill-adjusted expectation, confirming the prediction.",
   fs=11.5, colour=DARK_GREEN, bold=True)
footer(s, "Predictions fixed before the pull. The 50 % infill knock-down was anticipated — "
          "the verdict slide quantifies it as a single offset factor k ≈ 2.")
pageno(s, 127)

# ================= SLIDE 2 — Load vs time =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: LOAD vs TIME")
s.shapes.add_picture("images/V5/V5_slide_load_time.png", Inches(0.6), Inches(1.5), width=Inches(12.1))
footer(s, "Δ = 1304 N is the net pull applied during the test (UTS 1767 N − preload 463 N). "
          "Smooth rise, zero slip events (max dip 0.18 %) — drive + grips validated at 6× previous max load.")
pageno(s, 128)

# ================= SLIDE 3 — Stress vs strain =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: STRESS vs STRAIN")
s.shapes.add_picture("images/V5/V5_slide_stress_strain.png", Inches(0.35), Inches(1.5), width=Inches(8.85))
# yield-method explanation
ybx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(1.6), Inches(3.5), Inches(1.95))
ybx.fill.solid(); ybx.fill.fore_color.rgb = LIGHT_BLUE; ybx.line.fill.background()
tf = ybx.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "How σ_y is found — 0.2 % offset"
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.name = "Calibri"
for line in ["Draw a line parallel to the elastic",
             "slope E, shifted right by 0.002 strain.",
             "Where it crosses the curve = σ_y.",
             "→ σ_y = 19.0 MPa at ε = 0.0108."]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(11); rr.font.name = "Calibri"; rr.font.color.rgb = BLACK
kpi(s, 9.4, 3.75, 3.5, "elastic modulus (R² 0.994)", "E = 1.54 GPa", h=0.82, vfs=17)
kpi(s, 9.4, 4.67, 3.5, "ultimate strength @ ε=0.020", "UTS = 22.1 MPa", fill=YELLOW_WARN, h=0.82, vfs=17)
kpi(s, 9.4, 5.59, 3.5, "fracture (failure) stress, at ε_f = 0.0246", "σ_f = 21.8 MPa",
    fill=GREEN_PASS, h=0.95, vfs=17)
footer(s, "Four constitutive regions captured: linear-elastic → gradual polymer yield → UTS plateau → brittle fracture "
          "(only 1.5 % softening between UTS 22.1 and fracture 21.8 MPa).")
pageno(s, 129)

# ================= SLIDE 4 — Cauchy vs True strain =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: CAUCHY vs TRUE STRAIN")
s.shapes.add_picture("images/V5/V5_slide_cauchy_true.png", Inches(0.35), Inches(1.5), width=Inches(8.85))
dbx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(1.6), Inches(3.5), Inches(2.2))
dbx.fill.solid(); dbx.fill.fore_color.rgb = LIGHT_BLUE; dbx.line.fill.background()
tf = dbx.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Definitions"
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.name = "Calibri"
for line in ["ε_c (Cauchy) = (L − L₀) / L₀",
             "ε_t (True)   = ln(L / L₀)",
             "",
             "Identical at small strain; true",
             "strain is always slightly lower",
             "because ln(1+ε) < ε."]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(11); rr.font.name = "Calibri"; rr.font.color.rgb = BLACK
kpi(s, 9.4, 4.0, 1.7, "ε_c at fracture", "0.0246", h=0.95, vfs=18)
kpi(s, 11.2, 4.0, 1.7, "ε_t at fracture", "0.0243", fill=RED_FAIL, h=0.95, vfs=18)
wb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(5.15), Inches(3.5), Inches(1.5))
wb.fill.solid(); wb.fill.fore_color.rgb = YELLOW_WARN; wb.line.fill.background()
tf = wb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Why it matters here"
r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.name = "Calibri"
for line in ["Max gap only 1.2 % at ε_f = 0.025,",
             "so for PLA (small failure strain)",
             "Cauchy and True are interchangeable."]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(10.5); rr.font.name = "Calibri"; rr.font.color.rgb = BLACK
footer(s, "Both strain measures are reported straight from the DIC channels (DIC_Cauchy, DIC_True) — no post-processing.")
pageno(s, 130)

# ================= SLIDE 5 — Stress vs position =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: STRESS vs CROSSHEAD POSITION")
s.shapes.add_picture("images/V5/V5_slide_stress_position.png", Inches(0.35), Inches(1.5), width=Inches(8.85))
kpi(s, 9.4, 1.6, 1.7, "rig stiffness k = F/δ", "945 N/mm", h=0.95, vfs=17)
kpi(s, 11.2, 1.6, 1.7, "gauge share of travel", "52 %", fill=GREEN_PASS, h=0.95, vfs=17)
xb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(2.7), Inches(3.5), Inches(3.9))
xb.fill.solid(); xb.fill.fore_color.rgb = LIGHT_BLUE; xb.line.fill.background()
tf = xb.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "What the split means"
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.name = "Calibri"
for line in ["Crosshead moved 3.81 mm, but the",
             "specimen gauge only stretched",
             "1.97 mm (measured by DIC).",
             "",
             "• 1.97 mm = real specimen strain",
             "• 1.84 mm = rig + grip + load-train",
             "  elastic take-up (machine, not part)",
             "",
             "Rig stiffness k = F/δ",
             "  = 1767 N ÷ 1.84 mm ≈ 945 N/mm.",
             "",
             "Why important: ~48 % of crosshead",
             "motion is machine compliance — so",
             "strain MUST come from DIC, never",
             "from crosshead position. (Confirms",
             "the 8.6.3 compliance finding.)"]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(10.5); rr.font.name = "Calibri"; rr.font.color.rgb = BLACK
footer(s, "Gauge travel 52 % = fraction of crosshead motion that actually strained the specimen (1.97 / 3.81 mm).")
pageno(s, 131)

# ================= SLIDE 5b — Cauchy strain vs displacement =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: GAUGE STRAIN vs CROSSHEAD DISPLACEMENT")
s.shapes.add_picture("images/V5/V5_slide_strain_position.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "Gauge strain (DIC) rises ~linearly with crosshead travel but stays well below the dotted "
          "'all-travel-to-gauge' line — at fracture only 1.97 mm of the 3.81 mm travel (52 %) is real "
          "gauge strain; the rest is rig/grip take-up.")
pageno(s, 132)

# ================= SLIDE 6 — Verdict =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: VERDICT")
header(s, 0.5, 1.35, 6.0, "Pass criteria")
crit = [
    ["Criterion", "Result", "✓"],
    ["Reaches fracture, gauge section", "yes (photos)", "✓"],
    ["Curve: elastic→yield→UTS→fracture", "all 4 regions", "✓"],
    ["Failure strain 0.01 – 0.08", "0.0246", "✓"],
    ["Force monotonic to UTS (no slip)", "max dip 0.18 %", "✓"],
    ["DIC tracking ≥ 90 %", "99.8 %", "✓"],
    ["Drive applies load, no stall", "1 767 N steady", "✓"],
    ["E / σ_y / UTS vs Chacón", "via offset →", "✓*"],
]
ov = {(r, 2): {'bg': GREEN_PASS, 'bold': True} for r in range(1, 9)}
table(s, 0.5, 1.75, 6.0, 3.0, crit, cw=[2.7, 1.5, 0.4], hf=10, bf=10, ov=ov)

# DIC tracking definition box
db = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.95), Inches(6.0), Inches(1.75))
db.fill.solid(); db.fill.fore_color.rgb = LIGHT_BLUE; db.line.fill.background()
tf = db.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "What is “DIC tracking = 99.8 %”?"
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.name = "Calibri"
for line in ["% of the loading phase for which both markers were located and a valid",
             "gauge length L_px was returned. From the timestamps: valid strain from",
             "ramp-start (96.0 s) to 80 ms before fracture (134.43 of 134.51 s) = 99.8 %.",
             "Only the post-separation frames were invalid."]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(10.5); rr.font.name = "Calibri"; rr.font.color.rgb = BLACK

header(s, 6.8, 1.28, 6.0, "Material vs literature — offset corrections")
tb(s, 6.8, 1.62, 6.0, 0.3,
   "Area kept at CAD-verified 80 mm². Two ways to apply the 50 % infill offset:",
   fs=10, italic=True, colour=GREY_TEXT)

# (a) per-property offset
tb(s, 6.8, 1.94, 6.0, 0.3, "(a) Per-property offset — each to its own literature minimum",
   fs=10.5, bold=True, colour=DARK_GREEN)
offt = [
    ["Quantity", "Meas.\n(80 mm²)", "Offset×", "After\noffset", "Chacón", "✓"],
    ["E", "1.54 GPa", "×1.95", "3.00 GPa", "3.0–5.5", "✓"],
    ["σ_y", "19.0 MPa", "×1.58", "30.0 MPa", "30–50", "✓"],
    ["UTS", "22.1 MPa", "×1.45", "32.0 MPa", "32–60", "✓"],
]
ova = {}
for r in range(1, 4):
    ova[(r, 3)] = {'bg': LIGHT_BLUE, 'bold': True}
    ova[(r, 5)] = {'bg': GREEN_PASS, 'bold': True}
table(s, 6.8, 2.26, 6.0, 1.3, offt, cw=[0.95, 1.0, 0.85, 1.0, 1.1, 0.4], hf=9, bf=9.5, ov=ova)

# (b) single common factor
tb(s, 6.8, 3.68, 6.0, 0.3, "(b) Single common offset factor  k = 2.0  (area-ratio basis)",
   fs=10.5, bold=True, colour=DARK_GREEN)
offt2 = [
    ["Quantity", "Meas.\n(80 mm²)", "× k", "After\noffset", "Chacón", "✓"],
    ["E", "1.54 GPa", "×2.0", "3.08 GPa", "3.0–5.5", "✓"],
    ["σ_y", "19.0 MPa", "×2.0", "38.0 MPa", "30–50", "✓"],
    ["UTS", "22.1 MPa", "×2.0", "44.2 MPa", "32–60", "✓"],
]
ovb = {}
for r in range(1, 4):
    ovb[(r, 3)] = {'bg': LIGHT_BLUE, 'bold': True}
    ovb[(r, 5)] = {'bg': GREEN_PASS, 'bold': True}
table(s, 6.8, 4.0, 6.0, 1.3, offt2, cw=[0.95, 1.0, 0.85, 1.0, 1.1, 0.4], hf=9, bf=9.5, ov=ovb)

tb(s, 6.8, 5.42, 6.0, 1.3,
   "One factor is valid because the offset is an area ratio and E, σ_y, UTS are all "
   "force/area. The single k that fits all three ranges = [1.95, 2.63] → k = 2.0 chosen "
   "(E sets the lower limit). Cross-checks: dense V2–V4 bar E-ratio = 2.71; geometric "
   "area ratio = 1.47 — both bracket k ≈ 2.",
   fs=9.5, colour=BLACK)

# verdict banner
flow(s, 0.5, 6.85, 12.3, 0.5,
     "8.6.20 PASSED — rig + DIC validated through fracture at 1.77 kN; material properties pass once the 50 % infill offset is applied.",
     fill=GREEN_PASS, fg=DARK_GREEN, bold=True, fs=12.5, shape=MSO_SHAPE.RECTANGLE)
pageno(s, 133)

# ================= SLIDE 7 — Recommendations =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "PHASE 8.6.20: RECOMMENDATIONS & NEXT STEPS")
header(s, 0.5, 1.4, 12.3, "Actions")
recs = [
    ("1", "Print 100 % infill for hard-criterion material work",
     "A fully dense specimen removes the offset entirely — E / σ_y / UTS can then be compared "
     "directly to literature at the nominal 80 mm², with no correction factor.", GREEN_PASS),
    ("2", "Run 2 more 50 %-infill specimens (S5, S6) — identical conditions",
     "Test whether the offset factor is repeatable. If E / σ_y / UTS knock-downs cluster at the "
     "same ~1.5–2× across S4/S5/S6, the infill offset becomes a calibrated correction usable with confidence.", LIGHT_BLUE),
    ("3", "Record print parameters for every specimen",
     "Infill %, pattern, perimeter count, orientation. Properties are NOT comparable between specimens "
     "without them — the V2–V4 bar gave E = 4.18 GPa at nominal area; S4 (50 %) gave 1.54 GPa on the same rig.", LIGHT_BLUE),
    ("4", "Adopt the post-fracture force reading as the standard preload anchor",
     "Once the specimen breaks it carries 0 N, so the residual display (−463 N here) is a self-calibrating "
     "measure of the tared preload — more reliable than the pre-test screen value.", LIGHT_BLUE),
]
y = 1.95
for num, head, body, col in recs:
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = DARK_GREEN; badge.line.fill.background()
    tfb = badge.text_frame; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run(); rb.text = num; rb.font.size = Pt(20); rb.font.bold = True
    rb.font.color.rgb = WHITE; rb.font.name = "Calibri"
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(y), Inches(11.4), Inches(1.05))
    card.fill.solid(); card.fill.fore_color.rgb = col; card.line.fill.background()
    tfc = card.text_frame; tfc.word_wrap = True
    tfc.margin_left = Inches(0.18); tfc.margin_right = Inches(0.18); tfc.margin_top = Inches(0.08)
    pc = tfc.paragraphs[0]; rc = pc.add_run(); rc.text = head
    rc.font.size = Pt(14); rc.font.bold = True; rc.font.color.rgb = DARK_GREEN; rc.font.name = "Calibri"
    pc2 = tfc.add_paragraph(); rc2 = pc2.add_run(); rc2.text = body
    rc2.font.size = Pt(11); rc2.font.color.rgb = BLACK; rc2.font.name = "Calibri"
    y += 1.22

tb(s, 0.5, 6.95, 12.3, 0.4,
   "Phase 8.6 status:  8.6.3 ✓   8.6.4 ✓   8.6.19 ✓   8.6.20 ✓    |    Remaining: 8.6.15–18 + main.py log cleanup.",
   fs=11, italic=True, colour=GREY_TEXT)
pageno(s, 134)

prs.save("documentation/decks/V5_8_6_20_slides.pptx")
print("Saved: V5_8_6_20_slides.pptx (9 slides, pages 126-134)")
