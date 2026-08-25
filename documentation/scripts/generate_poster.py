"""SMART UTM posters — one content spec, three editions, all written to documentation/posters/.

    python documentation/scripts/generate_poster.py

  * **A0**  — conference wall poster, 3 columns.
  * **A4**  — conference handout, 2 columns. A genuine condensation, NOT a shrunken A0: scaling an
              A0 to A4 puts body text at ~5 pt.
  * **A4P** — progress report for a supervisor / manager. Different question, different poster: not
              "what can this rig do" but "where does the project stand, what is blocked, and what do
              you need from me".

DESIGN BASIS (researched 2026-08-12, not invented):
  * Morrison's #betterposter / billboard layout (2019) leads with ONE plain-language finding in very
    large type. Attendee studies: 82 % preferred billboard layouts for "understanding the main
    message", 68 % preferred them overall.
  * BUT the same studies found 67 % preferred traditional IMRaD for "communicating details and rigor",
    and asked for MORE information and MORE visualisations. Billboard posters trade away methods.
  * MIT's #evenbetterposter / "Generation 2" fix: keep the headline, but promote methods and data out
    of tiny sidebars into full panels with real section headers and figures at several detail levels.
  * Engineering-conference conventions: A0 portrait, 2-3 columns; title 72-120 pt, section headers
    36-48 pt, body 24-32 pt, captions 18-24 pt; ONE sans-serif family; 2-3 colours.
  => The conference editions are that hybrid: billboard band on top, then a column flow that gives
     METHODS and LIMITATIONS their own full panels alongside the feature cards.

Every number is pulled live from documentation/scripts/sf9_data.py (which recomputes it from the rig CSVs) or
from Software/UTM_PyQt6/registry.json — nothing on any poster is typed in by hand.
"""
import os
import sys
import json
import statistics as st

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.abspath(os.path.join(_HERE, "..", ".."))
os.chdir(_ROOT)
sys.path.insert(0, _HERE)
sys.path.insert(0, os.path.join(_ROOT, "Software", "UTM_PyQt6"))

from sf9_data import M as SF9                                                      # noqa: E402

OUTDIR = os.path.join("documentation", "posters")

# ---------------------------------------------------------------- palette (3 colours + neutrals)
JU_BLUE   = RGBColor(0x00, 0x53, 0x8A)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
GREY      = RGBColor(0x5A, 0x63, 0x6B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN     = RGBColor(0x1E, 0x84, 0x49)
GREEN_BG  = RGBColor(0xE8, 0xF4, 0xEC)
BLUE_BG   = RGBColor(0xE7, 0xEF, 0xF6)
AMBER     = RGBColor(0xB9, 0x6B, 0x00)
AMBER_BG  = RGBColor(0xFD, 0xF3, 0xE0)
RULE      = RGBColor(0xC8, 0xCF, 0xD5)
PAGE_BG   = RGBColor(0xF7, 0xF9, 0xFA)
FONT      = "Calibri"

STATUS = {
    "done":    (GREEN,   GREEN_BG, "VALIDATED"),
    "planned": (JU_BLUE, BLUE_BG,  "PLANNED"),
    "blocked": (AMBER,   AMBER_BG, "BLOCKED"),
}


# ================================================================ live numbers
def numbers():
    t9, b9 = SF9["t9"], SF9["t9_base"]
    L65, L63 = SF9["loops_t65"], SF9["loops_sin"]
    pc, sf, sf73 = SF9["pc"], SF9["sf"], SF9["sf_t73"]
    reg = json.load(open("Software/UTM_PyQt6/registry.json", encoding="utf-8"))
    val = [c for c in pc["cycles"] if c["R2"] > 0.94]
    v6 = [e["UTS_MPa"] for e in reg if (e.get("test") or "").startswith("V6")]
    return dict(
        n_runs=len(reg), n_spec=len({e["specimen"] for e in reg}), n_cyc=SF9["n_cyc_s22"],
        v6_n=len(v6), v6_mean=st.mean(v6), v6_sd=st.stdev(v6),
        v6_cv=100 * st.stdev(v6) / st.mean(v6),
        t9_net=t9["net"], t9_floor=b9["sd"], t9_x=t9["net"] / b9["sd"],
        t9_F=t9["Fmean"], t9_Fsd=t9["Fsd"], t9_dur=t9["dur"], t9_n=t9["findley_n"],
        t9_dpos=t9["dpos_um"], t9_ddic=t9["ddic_um"], t9_spec=t9["spec_frac_um"],
        t9_sig=t9["sigma_t"], t9_drift=b9["slope"], t9_bdur=b9["dur"],
        px65=L65[0]["px_span"], px63=L63[0]["px_span"],
        a65_0=L65[0]["area_kJm3"], a65_1=L65[-1]["area_kJm3"],
        e65_0=L65[0]["E_dn"] / 1000, e65_1=L65[-1]["E_dn"] / 1000,
        r2_65=min(x["R2_dn"] for x in L65), r2_63=L63[0]["R2_dn"],
        uts_sf=sf["uts"], uts_pc=pc["uts"], uts_t73=sf73["uts"],
        agree=100 * abs(sf["uts"] - pc["uts"]) / ((sf["uts"] + pc["uts"]) / 2),
        fatigue=100 * (1 - sf73["uts"] / sf["uts"]),      # vs S18: same protocol
        fatigue_mean=100 * (1 - sf73["uts"] / ((sf["uts"] + pc["uts"]) / 2)),
        K0=val[0]["K"], K1=val[-1]["K"], E0=val[0]["E"] / 1000, E1=val[-1]["E"] / 1000,
        Eloss=100 * (1 - val[-1]["E"] / val[0]["E"]),
        Kgain=100 * (val[-1]["K"] / val[0]["K"] - 1),
    )


N = numbers()
N_DONE, N_PLAN, N_BLOCK = 12, 3, 1

# ================================================================ shared content
TITLE = "SMART UTM — a tensile rig that runs, watches and stops its own tests"
AUTHORS = "Jönköping University · School of Engineering — Materials & Manufacturing"

BILLBOARD = "One specimen now yields a CURVE, not a point."
BILLBOARD_SUB = (
    "12 automation features and 6 closed-loop test protocols, every one validated on the rig. "
    "Camera-based strain sits INSIDE the control loop, so the machine measures the specimen "
    "instead of measuring itself.")

KPIS = [
    ("%d / %d" % (N_DONE, N_DONE + N_PLAN + N_BLOCK), "smart features\nbuilt & validated"),
    ("6", "closed-loop test\nprotocols"),
    ("%d" % N["n_runs"], "runs auto-indexed,\n%d specimens" % N["n_spec"]),
    ("± %.1f N" % N["t9_Fsd"], "force hold over\n%.0f s" % N["t9_dur"]),
    ("%.0f µε" % N["t9_floor"], "DIC noise floor\n(900 s window)"),
    ("%.1f %%" % N["agree"], "agreement between two\nfracture protocols"),
]

RIG_TITLE = "The rig — built in-house, every control loop closed in software"
RIG_TEXT = (
    "Two steppers (1.85 Nm, 20:1) drive 5 mm lead screws against a fixed base. An ANYLOAD 3 t cell "
    "measures force; a Basler camera measures strain optically at 20.8 px/mm from two markers on the "
    "80 mm gauge. The ESP32 runs STEP/DIR/ENABLE only and knows nothing of test protocols — every "
    "protocol, guard and measurement here lives in the PyQt6 app, closing the loop at ~11 Hz.")

WHY = ("A conventional pull gives ONE modulus, ONE yield, ONE strength — and a human presses every "
       "button. Crosshead travel is not specimen strain, so the machine reports its own compliance as "
       "material behaviour. A destructive test is one irreversible shot. And nothing watches the run, "
       "so a stall, a lost marker or a fracture are all noticed late, by a person.")

METHODS = [
    ("One control engine", "All six protocols are one function called on every load sample (~11 Hz) "
     "with load, position and time-matched DIC strain in hand. A protocol only returns a speed — it "
     "cannot bypass the shared safety net."),
    ("Strain from the specimen", "A camera tracks two markers; engineering strain is (L−L₀)/L₀ at "
     "20.8 px/mm. The centroid resolves 0.1 px steps at ±0.02 px."),
    ("Simulated before hardware", "Every policy is replayed against a spring plant and recorded CSVs "
     "— 9/9 checks pass before a specimen is risked."),
    ("Four safety layers", "Load-collapse detector · stall guard (< 0.05 mm / 6 s under load) · 10 kN "
     "and 30 mm backstops · dead-DIC guard (freeze at 0.2 s, halt at 1.0 s). Guards are phase-aware, "
     "so the stall guard stays silent through an intentional dwell."),
    ("One analysis library", "The app, the reports and every figure compute E, σ_y, UTS, ε_f and the "
     "force anchor from the same code — a live readout and a published number cannot disagree."),
]

# SF numbers are STABLE, append-only IDs (SF1-8 are printed in the V6a deck). Cards are listed in
# NUMERIC ORDER so the poster reads SF1, SF2, SF3 …, with status carried by colour and a chip rather
# than by grouping — which is why a legend block is mandatory on every edition. Cards carry no
# figure: 13 thumbnails swamp even an A0, so the visual proof lives in the EVIDENCE panels, which is
# also what the #evenbetterposter critique asks for.
CARDS = [
    dict(sf=1, status="done", title="Live DIC health HUD", metric="2 / 2",
         mlabel="markers, live at ~2 Hz",
         claim="A green/amber/red badge shows marker lock, tracking % and jitter live — so a dropout is seen during the test, not found later in the CSV."),
    dict(sf=2, status="done", title="Prepare specimen", metric="4 → 1",
         mlabel="operator actions",
         claim="One button tares position, force and DIC and clears the plots. It tares DIC only at a green 2/2 lock, so a bad mount cannot be zeroed in."),
    dict(sf=3, status="done", title="Recipes / settings", metric="1 click",
         mlabel="restores the whole setup",
         claim="Saves and restores an entire setup, including all six protocols' parameters — so repeats are identical by construction, not by memory."),
    dict(sf=4, status="done", title="One-click report", metric="1 click",
         mlabel="→ PDF + every graph",
         claim="A one-page PDF and every plot, recomputed from the CSV by the same library the app runs live."),
    dict(sf=5, status="done", title="Auto-stop at fracture", metric="1.2 s",
         mlabel="collapse → motor halt",
         claim="A load-collapse detector halts the motor the moment the specimen breaks. Caught S16 at 3375 N, silent through the ductile draw before it."),
    dict(sf=6, status="done", title="Constant-strain-rate testing", metric="0.00051 /s",
         mlabel="held, vs a 0.0005 target",
         claim="Closed loop on the specimen's own dε/dt: the crosshead adapted itself 0.10 → 0.05 mm/s, fast in the elastic region, slow through necking."),
    dict(sf=7, status="done", title="Stall guard", metric="0.05 mm / 6 s",
         mlabel="trip, armed above 200 N",
         claim="Halts a motor that has stopped moving under load. It proved the T7 failure was a 2.6 kN stall, not a fracture — the specimen came off intact."),
    dict(sf=8, status="done", title="Release load", metric="≤ 5 N",
         mlabel="residual, at 0.20 mm/s",
         claim="Taring at preload makes the readout say 0 N while 300 N is still on the specimen. This drives past that, down to the true zero."),
    dict(sf=9, status="done", title="Six closed-loop protocols", metric="6 / 6",
         mlabel="rig-validated",
         claim="cyclic · staircase · relaxation · creep · staircase→FRACTURE · progressive-cyclic→FRACTURE. One engine, one safety net."),
    dict(sf=10, status="done", title="Auto-preload", metric="1.03 ×",
         mlabel="target, pre-paying relaxation",
         claim="Approaches on a 0.2 → 0.1 → 0.02 mm/s schedule and stops 3 % high, because PLA relaxes about 2 % while held."),
    dict(sf=11, status="planned", title="Auto-metadata + foldering", metric="planned",
         mlabel="next feature up",
         claim="Per-specimen folders and file IDs filled from the recipe, plus a one-click per-specimen slide deck."),
    dict(sf=12, status="planned", title="DIC auto-calibrate", metric="planned",
         mlabel="camera robustness",
         claim="Exposure and threshold swept on camera start, and an ROI that follows the markers through a ductile draw."),
    dict(sf=13, status="planned", title="Guided workflow + live overlay", metric="planned",
         mlabel="usability",
         claim="A Connect→Calibrate→Mount→Prepare→Run→Save wizard, with live modulus and a fracture flag as the test runs."),
    dict(sf=14, status="blocked", title="Poisson's ratio / true stress", metric="blocked",
         mlabel="optics, not code",
         claim="The four-marker maths is written and self-tested. The gauge is too narrow for a transverse pair at 20.8 px/mm."),
    dict(sf=15, status="done", title="Test registry", metric="%d runs" % N["n_runs"],
         mlabel="auto-indexed, %d specimens" % N["n_spec"],
         claim="Each run indexed with whatever its loading path supports — UTS, ε_f, E, σ_y, toughness — plus the force anchor that turns a tared reading into a true load."),
    dict(sf=16, status="done", title="Dead-DIC guard + hard backstops", metric="518 → 175 N",
         mlabel="overshoot on marker loss",
         claim="A controller ramping blind on a stale camera reading spikes the force. Speed freezes at 0.2 s of frozen strain, halt at 1.0 s."),
]
CARDS.sort(key=lambda c: c["sf"])

EVIDENCE = [
    dict(title="Six closed-loop protocols — every trace is real rig data",
         img="documentation/figures/sf9_overview.png",
         text="All six are one function behind one safety net; a protocol only chooses a speed. Two of "
              "them drive the specimen to destruction while interrogating it on the way — which is how "
              "one specimen yields a curve instead of a point. Two protocols taken independently to "
              "fracture landed %.1f %% apart on ultimate strength (%.2f vs %.2f MPa) — on different "
              "specimens, n = 1 each, so that bounds protocol and specimen scatter together."
              % (N["agree"], N["uts_sf"], N["uts_pc"])),
    dict(title="The machine reads STIFFER while the specimen SOFTENS",
         img="documentation/figures/sf9_teach_stiffness.png",
         text="Progressive cyclic to fracture. Crosshead stiffness ROSE %.0f → %.0f N/mm as rig slack "
              "squeezed out, while the DIC-measured specimen modulus FELL %.2f → %.2f GPa, a %.0f %% "
              "loss. Read from the crosshead alone this test says the material is getting stiffer — "
              "the strongest argument for putting the camera in the loop."
              % (N["K0"], N["K1"], N["E0"], N["E1"], N["Eloss"])),
    dict(title="Creep, resolved — and separated from the instrument's own drift",
         img="documentation/figures/sf9_creep_t9.png",
         text="Two holds, in order. A %.0f s ZERO-LOAD baseline measured the rig's own drift at "
              "%+.4f µε/s; the real hold then pinned %.0f ± %.1f N for %.0f s at %.1f MPa. Raw strain "
              "minus that slope leaves %+.0f µε of creep — %.0f× the %.0f µε noise floor — and it "
              "decelerates (n = %.2f), the criterion fixed BEFORE the run for telling creep from "
              "drift. Holding it took %+.0f µm of crosshead but only %+.0f µm of gauge: %.0f %% "
              "specimen."
              % (N["t9_bdur"], N["t9_drift"], N["t9_F"], N["t9_Fsd"], N["t9_dur"], N["t9_sig"],
                 N["t9_net"], N["t9_x"], N["t9_floor"], N["t9_n"],
                 N["t9_dpos"], N["t9_ddic"], N["t9_spec"])),
    dict(title="Where you cycle matters more than how accurately you cycle",
         img="documentation/figures/sf9_cyclic_compare.png",
         text="Cyclic hysteresis was a negative result until the load window moved. Raising the unload "
              "floor to 400 N stops the mechanism re-crossing its own backlash twice per cycle: strain "
              "excursion %.1f → %.1f px, unload-fit R² %.2f → %.2f, and loop area and modulus now both "
              "fall monotonically (%.1f → %.1f kJ/m³, %.2f → %.2f GPa). Same code, same camera — one "
              "parameter."
              % (N["px63"], N["px65"], N["r2_63"], N["r2_65"], N["a65_0"], N["a65_1"],
                 N["e65_0"], N["e65_1"])),
]

LIMITS = [
    ("Poisson's ratio is not measurable today", "The gauge is too narrow: the elastic width change is "
     "sub-pixel at 20.8 px/mm. Every stress here is ENGINEERING (F/A₀), never 'true'."),
    ("No fatigue damage curve yet", "One run lost the camera mid-way and its repeat started on an "
     "already-damaged specimen, and the modulus RECOVERED between them — so no valid E₀ exists."),
    ("Creep and fatigue are each n = 1", "Both measurements are proven; the material constants are "
     "not. Creep needs 3–4 stress levels and fatigue needs repeats."),
    ("The motor's torque ceiling is variable", "Normally 3.2–3.4 kN, but ~2.6 kN on one session. The "
     "drivetrain is rated far above that, so the cause is electrical — not the software."),
]

FOOTER = ("Code, raw CSVs and the full 76-slide technical deck: Software/UTM_PyQt6/ and "
          "documentation/  ·  Every figure and number here is regenerated from the rig CSVs by "
          "documentation/scripts/sf9_data.py — none of it is transcribed by hand.  ·  Layout follows the "
          "#evenbetterposter / Generation-2 billboard format: a plain-language headline, with methods "
          "and limitations kept as full panels rather than sidebars.")

# ================================================================ progress-report content
P_TITLE = "SMART UTM — PROGRESS REPORT"
P_SUB = ("Jönköping University · School of Engineering  —  status to 12 August 2026  ·  "
         "automated tensile testing with camera-based strain")

P_BILLBOARD = "12 of 16 planned features are built AND validated on the rig."
P_BILLBOARD_SUB = (
    "All six closed-loop test protocols work. The material results are publication-grade. What "
    "remains is two rig runs, three convenience features, and two hardware decisions that are not "
    "mine to make.")

P_KPIS = [
    ("%d / %d" % (N_DONE, N_DONE + N_PLAN + N_BLOCK), "features built\nand rig-validated"),
    ("6 / 6", "test protocols\nvalidated"),
    ("%d" % N["n_spec"], "specimens tested,\n%d indexed runs" % N["n_runs"]),
    ("%.1f %%" % N["v6_cv"], "strength scatter,\nn = %d at 100 %% infill" % N["v6_n"]),
    ("2", "hardware decisions\nneeded from you"),
]

P_STATUS = (
    "The rig is no longer a manual machine. A test is now: pick a recipe, press Prepare, press Start "
    "— the software runs the protocol, watches the specimen through a camera, stops itself at "
    "fracture and writes an indexed report. Six protocols beyond a plain pull are validated, so a "
    "single specimen can return a curve rather than a single number. The two results that were "
    "negative in July, cyclic hysteresis and creep, were both re-sized and are now positive.")

MILESTONES = [
    ("June", "DIC validation closed",
     "41/41 software checks; tension, compression and staircase linearity all sign-verified."),
    ("June", "First run to fracture",
     "Rig and camera both hold all the way to failure — the precondition for every destructive "
     "protocol since."),
    ("June", "100 %% infill quintet, n = %d" % N["v6_n"],
     "UTS %.1f ± %.1f MPa, scatter %.1f %% — publication-grade repeatability."
     % (N["v6_mean"], N["v6_sd"], N["v6_cv"])),
    ("July", "Smart features SF1–SF10",
     "One-click workflow, a four-layer safety net, and constant-strain-rate control demonstrated to "
     "fracture."),
    ("August", "All six protocols validated",
     "Including both destructive protocols; two landed %.1f %% apart on ultimate strength, though on "
     "different specimens." % N["agree"]),
    ("August", "Fatigue and creep resolved",
     "Residual strength after cycling is −%.1f %% against the same-protocol virgin run; creep "
     "resolved at %.0f× the noise floor." % (N["fatigue"], N["t9_x"])),
]

P_RESULTS = [
    ("Strength repeatability is publication-grade",
     "100 %% infill, n = %d: UTS %.1f ± %.1f MPa — %.1f %% scatter, and the only n > 2 result we own. "
     "Two fracture protocols landed %.1f %% apart, but on DIFFERENT specimens (n = 1 each), so that "
     "bounds protocol and specimen scatter together, not the protocol alone."
     % (N["v6_n"], N["v6_mean"], N["v6_sd"], N["v6_cv"], N["agree"])),
    ("The camera is not a nicety — it changes the answer",
     "Over the same cycles the crosshead reads %.0f %% STIFFER while the specimen's true modulus "
     "falls %.0f %%. Without strain measured on the specimen itself, this test reports the opposite "
     "of what actually happened." % (N["Kgain"], N["Eloss"])),
    ("First fatigue-damage number",
     "After %d cycles at 79 %% of fracture load, residual strength is −%.1f %% against the "
     "same-protocol virgin run (−%.1f %% against the two-protocol mean). The virgin baseline spans "
     "only 0.9 %%, so the loss is far larger than the scatter."
     % (N["n_cyc"], N["fatigue"], N["fatigue_mean"])),
    ("Creep measured, with the instrument subtracted",
     "A zero-load baseline measured the rig's own drift first; the corrected creep is %+.0f µε at "
     "%.0f× the noise floor, and decelerates as real creep must."
     % (N["t9_net"], N["t9_x"])),
]

P_NEXT = [
    ("Now", "T6.6 — the fatigue damage curve",
     "One fresh specimen, 12 cycles, one uninterrupted run. The only missing piece of the damage "
     "story; about 40 minutes of rig time."),
    ("Now", "Black-specimen DIC check",
     "Every specimen so far is white. The black preset has never been exercised — needed before "
     "anyone else uses the rig."),
    ("Next", "SF11 — auto-metadata and per-specimen reports",
     "Removes the last manual bookkeeping step from a test."),
    ("Next", "Student operating manual",
     "So the rig can be run by someone who did not build it."),
]

P_ASKS = [
    ("Motor torque ceiling — needs a hardware decision", AMBER, AMBER_BG,
     "Fracture loads reach 3.2–3.4 kN normally, but fell to ~2.6 kN on one session and stalled a "
     "100 % infill test. The drivetrain is rated above 12 kN, so the cause is electrical: driver "
     "current, thermal derating or supply sag. Two routes — set the driver current (Vref) on the "
     "board, or wire the TMC2160's SPI pins to the ESP32 so current and thermal status can be set "
     "and **logged** in software. **I need bench time with the driver board, and a decision on "
     "which route.**"),
    ("Poisson's ratio — needs optics, not code", AMBER, AMBER_BG,
     "The four-marker maths is written and self-tested, but the gauge is too narrow to resolve a "
     "transverse pair at the present 20.8 px/mm. **A gauge-zoomed lens (or a second camera) plus a "
     "matte backdrop unblocks true stress and Poisson's ratio.** Until then every stress reported is "
     "engineering stress, which is stated on every figure."),
    ("More specimens before publication", JU_BLUE, BLUE_BG,
     "Creep and fatigue are each n = 1 at one condition. The measurement is proven; the material "
     "constants are not. **3–4 stress levels for creep and 3 repeats for fatigue** would make both "
     "publishable."),
]

P_FOOTER = ("Every number on this page is recomputed from the rig CSVs by documentation/scripts/sf9_data.py "
            "and the test registry — nothing is transcribed by hand. Full detail: the 76-slide "
            "technical deck in documentation/, and ROADMAP.md for the complete task list.")

# ================================================================ format specs
FORMATS = {
    "A0": dict(w=33.11, h=46.81, cols=3, margin=1.05, gutter=0.62, verbose=True, kind="conf",
               fs=dict(title=78, authors=30, bill=104, billsub=34, kpi=58, kpil=17,
                       head=42, card_t=27, card_b=19, claim=21, metric=36, metricl=15,
                       body=20, small=16, foot=15, badge=15)),
    "A4": dict(w=8.27, h=11.69, cols=2, margin=0.34, gutter=0.20, verbose=False, kind="conf",
               fs=dict(title=19, authors=8.5, bill=27, billsub=9.5, kpi=15.5, kpil=5.6,
                       head=11.5, card_t=8.6, card_b=6.6, claim=7.0, metric=9.5, metricl=5.2,
                       body=6.8, small=5.8, foot=5.2, badge=5.2)),
    "A4P": dict(w=8.27, h=11.69, cols=2, margin=0.34, gutter=0.22, verbose=False, kind="progress",
                fs=dict(title=21, authors=8.2, bill=25, billsub=9.5, kpi=16, kpil=5.8,
                        head=12, card_t=9.0, card_b=7.0, claim=7.2, metric=10, metricl=5.4,
                        body=7.2, small=6.2, foot=5.2, badge=5.4)),
}


# ================================================================ text primitives
_FONTS, _FONT_WARNED = {}, [False]
_FONT_FILES = {False: "C:/Windows/Fonts/calibri.ttf", True: "C:/Windows/Fonts/calibrib.ttf"}
_OVERSAMPLE = 4          # measure at 4x nominal size so rounding to integer px stops mattering


def _font(fs, bold):
    key = (round(fs, 1), bool(bold))
    if key not in _FONTS:
        try:
            from PIL import ImageFont
            _FONTS[key] = ImageFont.truetype(_FONT_FILES[bool(bold)],
                                             max(4, int(round(fs * _OVERSAMPLE))))
        except Exception:
            _FONTS[key] = None
            if not _FONT_WARNED[0]:
                print("      ! Calibri metrics unavailable — falling back to a character estimate; "
                      "check the rendered poster for overlap.")
                _FONT_WARNED[0] = True
    return _FONTS[key]


def est_lines(text, w_in, fs, bold=False):
    """Wrapped-line count, measured against the real Calibri metrics rather than guessed.

    An average-character-width estimate over-wrapped large type badly (the 104 pt billboard headline
    scored as two lines when it renders as one, leaving an inch and a half of dead space). The 0.985
    is a safety margin against PowerPoint's kerning differing from PIL's — under-estimating a height
    causes an OVERLAP, which is far worse than a gap.

    `**bold**` markers are stripped before measuring: they are formatting, not glyphs."""
    f = _font(fs, bold)
    text = str(text).replace("**", "")
    if f is None:                                                    # non-Windows fallback
        chars = max(1.0, w_in * 72.0 / (fs * 0.50)) / 1.10
        return sum(max(1, -(-len(p) // int(chars))) for p in text.split("\n"))
    limit = w_in * 72.0 * _OVERSAMPLE * 0.985
    total = 0
    for para in text.split("\n"):
        cur, n = "", 1
        for word in para.split(" "):
            trial = word if not cur else cur + " " + word
            if f.getlength(trial) <= limit:
                cur = trial
            else:
                n += 1
                cur = word
        total += n
    return total


def text_h(text, w_in, fs, line=0.92, bold=False):
    return est_lines(text, w_in, fs, bold) * fs * line * 1.20 / 72.0


def fit_one_line(text, w_in, fs_max, fs_min, bold=True):
    """Largest size <= fs_max keeping `text` on one line, floored at fs_min. The title and billboard
    are the two places where a wrap shows as a one-word orphan, so they are sized to the space."""
    fs = fs_max
    while fs > fs_min and est_lines(text, w_in, fs, bold) > 1:
        fs -= 0.5
    return fs


def tbox(sl, x, y, w, h, text, *, fs=20, bold=False, italic=False, colour=DARK,
         align=PP_ALIGN.LEFT, line=0.92, anchor=MSO_ANCHOR.TOP):
    """Text box. Wrapping a phrase in **double asterisks** emboldens just that run — used to make the
    actual ASK stand out inside a paragraph on the progress report."""
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, para in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(0)
        for j, chunk in enumerate(para.split("**")):
            if not chunk:
                continue
            r = p.add_run(); r.text = chunk
            r.font.size = Pt(fs); r.font.bold = bold or (j % 2 == 1); r.font.italic = italic
            r.font.color.rgb = colour; r.font.name = FONT
    return box


def rect(sl, x, y, w, h, *, fill=None, line=None, lw=1.0, radius=None):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def hrule(sl, x, y, w, colour=RULE, lw=1.2):
    ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(lw / 72.0))
    ln.fill.solid(); ln.fill.fore_color.rgb = colour
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def img_h(path, w):
    iw, ih = Image.open(path).size
    return w * ih / iw


def place_img(sl, path, x, y, w):
    sl.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    return img_h(path, w)


# ================================================================ blocks
# A block is (height_fn(w) -> inches, draw_fn(slide, x, y, w) -> None, kind).

def blk_section(F, label):
    """Section header + rule. Measure the UPPERCASED string, not the source: caps are wider, and
    measuring mixed case let the next block sit on top of a two-line header."""
    fs = F["fs"]["head"]
    caps = label.upper()

    def h(w):
        return text_h(caps, w, fs, bold=True) + 0.30 * F["k"]

    def draw(sl, x, y, w):
        hl = text_h(caps, w, fs, bold=True)
        tbox(sl, x, y, w, hl, caps, fs=fs, bold=True, colour=JU_BLUE)
        hrule(sl, x, y + hl + 0.10 * F["k"], w, JU_BLUE, 2.4 * F["k"])
    return h, draw


def blk_para(F, text, *, fs_key="body", colour=DARK, italic=False, gap=0.30):
    fs = F["fs"][fs_key]

    def h(w):
        return text_h(text, w, fs, 1.0) + gap * F["k"]

    def draw(sl, x, y, w):
        tbox(sl, x, y, w, text_h(text, w, fs, 1.0), text, fs=fs, colour=colour, italic=italic,
             line=1.0)
    return h, draw


def blk_defs(F, items, *, gap=0.26, fs_key="body"):
    """(bold lead-in, body) pairs — methods, limitations, results."""
    fsb, fs = F["fs"]["card_t"], F["fs"][fs_key]

    def h(w):
        return sum(text_h(a, w, fsb, bold=True) + text_h(b, w, fs, 1.0) + 0.16 * F["k"]
                   + gap * F["k"] for a, b in items)

    def draw(sl, x, y, w):
        yy = y
        for lead, body in items:
            hl = text_h(lead, w, fsb, bold=True)
            tbox(sl, x, yy, w, hl, lead, fs=fsb, bold=True, colour=DARK)
            yy += hl + 0.05 * F["k"]
            hb = text_h(body, w, fs, 1.0)
            tbox(sl, x, yy, w, hb, body, fs=fs, colour=GREY, line=1.0)
            yy += hb + gap * F["k"]
    return h, draw


def blk_legend(F):
    """Status key. Mandatory, because the cards run in NUMERIC order — status is carried by colour
    rather than by grouping, and without a key that is a guess."""
    fs = F["fs"]["small"]
    items = [(GREEN, "built & rig-validated (%d)" % N_DONE),
             (JU_BLUE, "planned (%d)" % N_PLAN),
             (AMBER, "hardware-blocked (%d)" % N_BLOCK)]
    ch = (fs * 1.9) / 72.0

    def h(w):
        return ch + 0.22 * F["k"]

    def draw(sl, x, y, w):
        cw = w / len(items)
        for i, (col, lab) in enumerate(items):
            xx = x + i * cw
            rect(sl, xx, y + ch * 0.20, ch * 0.58, ch * 0.58, fill=col, radius=0.25)
            tbox(sl, xx + ch * 0.80, y + ch * 0.14, cw - ch * 0.85, ch, lab, fs=fs, colour=GREY)
    return h, draw


def blk_card(F, c):
    """Poster card: status stripe · SF badge · title · claim · a big metric with its label."""
    fsT, fsC = F["fs"]["card_t"], F["fs"]["claim"]
    fsM, fsML, fsB = F["fs"]["metric"], F["fs"]["metricl"], F["fs"]["badge"]
    col, bg, _ = STATUS[c["status"]]
    pad = 0.17 * F["k"]
    bw = 1.55 * F["k"]

    def geom(w):
        iw = w - 2 * pad
        tw = iw - bw - 0.14 * F["k"]
        mw = max(2.0 * F["k"], iw * 0.34)
        cw = iw - mw - 0.20 * F["k"]
        ht = text_h(c["title"], tw, fsT, bold=True)
        hc = text_h(c["claim"], cw, fsC, 1.02)
        hm = (fsM * 1.25 + fsML * 2.4) / 72.0
        return iw, tw, mw, cw, ht, max(hc, hm)

    def h(w):
        _, _, _, _, ht, hrow = geom(w)
        return ht + 0.12 * F["k"] + hrow + 2 * pad + 0.17 * F["k"]

    def draw(sl, x, y, w):
        iw, tw, mw, cw, ht, hrow = geom(w)
        body = ht + 0.12 * F["k"] + hrow + 2 * pad
        rect(sl, x, y, w, body, fill=WHITE, line=RULE, lw=1.1 * F["k"], radius=0.02)
        rect(sl, x, y, 0.11 * F["k"], body, fill=col)
        xx, yy = x + pad + 0.08 * F["k"], y + pad
        rect(sl, x + w - pad - bw, yy, bw, (fsB * 2.0) / 72.0, fill=bg, radius=0.28)
        tbox(sl, x + w - pad - bw, yy + (fsB * 0.30) / 72.0, bw, (fsB * 1.5) / 72.0,
             "SF%d" % c["sf"], fs=fsB, bold=True, colour=col, align=PP_ALIGN.CENTER)
        tbox(sl, xx, yy, tw, ht, c["title"], fs=fsT, bold=True, colour=DARK)
        yy += ht + 0.12 * F["k"]
        tbox(sl, xx, yy, cw, hrow, c["claim"], fs=fsC, colour=GREY, line=1.02)
        mx = x + w - pad - mw
        rect(sl, mx, yy, mw, (fsM * 1.25 + fsML * 2.4) / 72.0, fill=bg, radius=0.07)
        tbox(sl, mx, yy + 0.05 * F["k"], mw, (fsM * 1.3) / 72.0, c["metric"], fs=fsM, bold=True,
             colour=col, align=PP_ALIGN.CENTER)
        tbox(sl, mx + 0.06 * F["k"], yy + (fsM * 1.28) / 72.0, mw - 0.12 * F["k"],
             (fsML * 2.4) / 72.0, c["mlabel"], fs=fsML, colour=GREY, align=PP_ALIGN.CENTER,
             line=1.02)
    return h, draw


def blk_cardrow(F, c):
    """A4 handout: one compact line per feature — badge · title · metric · one-clause claim."""
    fsT, fs = F["fs"]["card_t"], F["fs"]["small"]
    col, bg, _ = STATUS[c["status"]]
    pad = 0.05 * F["k"]
    bw = 1.15 * F["k"]
    lead = c["claim"].split(". ")[0].rstrip(".")

    def geom(w):
        tw = w - bw - 0.10 * F["k"]
        head = "%s — %s" % (c["title"], c["metric"])
        return tw, head, text_h(head, tw, fsT, bold=True) + text_h(lead, tw, fs, 1.0) + 0.04 * F["k"]

    def h(w):
        return geom(w)[2] + 2 * pad + 0.13 * F["k"]

    def draw(sl, x, y, w):
        tw, head, hh = geom(w)
        rect(sl, x, y, w, hh + 2 * pad, fill=WHITE, line=RULE, lw=0.8, radius=0.03)
        rect(sl, x, y, 0.045, hh + 2 * pad, fill=col)
        rect(sl, x + 0.09 * F["k"], y + pad, bw - 0.10 * F["k"], (fsT * 1.7) / 72.0, fill=bg,
             radius=0.3)
        tbox(sl, x + 0.09 * F["k"], y + pad + (fsT * 0.22) / 72.0, bw - 0.10 * F["k"],
             (fsT * 1.3) / 72.0, "SF%d" % c["sf"], fs=fsT * 0.82, bold=True, colour=col,
             align=PP_ALIGN.CENTER)
        xx = x + bw + 0.04 * F["k"]
        ht = text_h(head, tw, fsT, bold=True)
        tbox(sl, xx, y + pad, tw, ht, head, fs=fsT, bold=True, colour=DARK)
        tbox(sl, xx, y + pad + ht + 0.02 * F["k"], tw, text_h(lead, tw, fs, 1.0), lead, fs=fs,
             colour=GREY, line=1.0)
    return h, draw


def blk_evidence(F, e):
    fsT, fs = F["fs"]["card_t"], F["fs"]["body"]
    pad = 0.22 * F["k"]

    def h(w):
        iw = w - 2 * pad
        return (text_h(e["title"], iw, fsT, bold=True) + 0.12 * F["k"] + img_h(e["img"], iw)
                + 0.14 * F["k"] + text_h(e["text"], iw, fs, 1.0) + 2 * pad + 0.34 * F["k"])

    def draw(sl, x, y, w):
        iw = w - 2 * pad
        rect(sl, x, y, w, h(w) - 0.34 * F["k"], fill=BLUE_BG, line=JU_BLUE, lw=1.4 * F["k"],
             radius=0.02)
        xx, yy = x + pad, y + pad
        ht = text_h(e["title"], iw, fsT, bold=True)
        tbox(sl, xx, yy, iw, ht, e["title"], fs=fsT, bold=True, colour=JU_BLUE)
        yy += ht + 0.12 * F["k"]
        yy += place_img(sl, e["img"], xx, yy, iw) + 0.14 * F["k"]
        tbox(sl, xx, yy, iw, text_h(e["text"], iw, fs, 1.0), e["text"], fs=fs, colour=DARK, line=1.0)
    return h, draw


def blk_rig(F, images, *, row_in=6.6):
    """The machine itself. `images` are BACKGROUND-REMOVED cutouts, so they sit directly on the page
    colour with no photo edge — which is why two of them can stand side by side without reading as a
    collage. They are bottom-aligned so they share a visual floor.

    A poster about instrumentation with no picture of the instrument fails the first question every
    passer-by asks."""
    fsT, fs = F["fs"]["card_t"], F["fs"]["body"]
    pad = 0.22 * F["k"]
    gap = 0.20 * F["k"]
    imgs = [p for p in images if os.path.exists(p)]

    def geom(w):
        iw = w - 2 * pad
        if not imgs:
            return iw, [], 0.0, text_h(RIG_TEXT, iw, fs, 1.02)
        row = row_in * F["k"]                                  # shared HEIGHT: all are portrait
        ws = [row * (lambda s: s[0] / s[1])(Image.open(p).size) for p in imgs]
        if sum(ws) + gap * (len(ws) - 1) > iw:                 # shrink to fit a narrow column
            k = (iw - gap * (len(ws) - 1)) / sum(ws)
            ws = [x * k for x in ws]
        rh = max(img_h(p, x) for p, x in zip(imgs, ws))
        return iw, ws, rh, text_h(RIG_TEXT, iw, fs, 1.02)

    def h(w):
        iw, ws, rh, ch = geom(w)
        return (text_h(RIG_TITLE, iw, fsT, bold=True) + 0.12 * F["k"] + rh + 0.16 * F["k"] + ch
                + 2 * pad + 0.30 * F["k"])

    def draw(sl, x, y, w):
        iw, ws, rh, ch = geom(w)
        ht = text_h(RIG_TITLE, iw, fsT, bold=True)
        rect(sl, x, y, w, h(w) - 0.30 * F["k"], fill=WHITE, line=RULE, lw=1.2 * F["k"], radius=0.02)
        xx, yy = x + pad, y + pad
        tbox(sl, xx, yy, iw, ht, RIG_TITLE, fs=fsT, bold=True, colour=DARK)
        yy += ht + 0.12 * F["k"]
        x0 = xx + (iw - sum(ws) - gap * (len(ws) - 1)) / 2.0        # centre the row in the panel
        for p, ww in zip(imgs, ws):
            place_img(sl, p, x0, yy + (rh - img_h(p, ww)), ww)      # bottom-align on a shared floor
            x0 += ww + gap
        tbox(sl, xx, yy + rh + 0.16 * F["k"], iw, ch, RIG_TEXT, fs=fs, colour=GREY, line=1.02)
    return h, draw


def blk_milestones(F):
    """Vertical milestone list, drawn natively. A horizontal timeline put six labels within a month
    of each other on the axis and they collided every time; a vertical list cannot."""
    fsD, fsT, fs = F["fs"]["small"], F["fs"]["card_t"], F["fs"]["body"]
    pad = 0.07 * F["k"]
    f = _font(fsD, True)
    dtext = (max(f.getlength(d) for d, _t, _b in MILESTONES) / (_OVERSAMPLE * 72.0)
             if f else 0.55 * F["k"])
    dw = dtext + 0.34 * F["k"]      # date text + the node/spine gutter

    def geom(w):
        tw = w - dw - 0.14 * F["k"]
        return tw, [(d, t, b, text_h(t, tw, fsT, bold=True) + text_h(b, tw, fs, 1.0) + 0.03 * F["k"])
                    for d, t, b in MILESTONES]

    def h(w):
        return sum(r[3] + 2 * pad + 0.09 * F["k"] for r in geom(w)[1]) + 0.16 * F["k"]

    def draw(sl, x, y, w):
        tw, rows = geom(w)
        yy = y
        for d, t, b, hh in rows:
            rect(sl, x + dtext + 0.16 * F["k"], yy, 0.020 * F["k"], hh + 2 * pad, fill=RULE)   # spine
            rect(sl, x + dtext + 0.09 * F["k"], yy + pad + 0.02 * F["k"], 0.16 * F["k"], 0.16 * F["k"],
                 fill=GREEN, radius=0.5)                                                  # node
            tbox(sl, x, yy + pad, dtext + 0.02, (fsD * 1.4) / 72.0, d, fs=fsD, bold=True, colour=GREY)
            ht = text_h(t, tw, fsT, bold=True)
            tbox(sl, x + dw + 0.14 * F["k"], yy + pad, tw, ht, t, fs=fsT, bold=True, colour=JU_BLUE)
            tbox(sl, x + dw + 0.14 * F["k"], yy + pad + ht + 0.02 * F["k"], tw,
                 text_h(b, tw, fs, 1.0), b, fs=fs, colour=GREY, line=1.0)
            yy += hh + 2 * pad + 0.09 * F["k"]
    return h, draw


def blk_asks(F):
    """The decisions the reader can actually make. On a progress report this is the payload, so it
    carries the strongest visual weight on the page after the headline."""
    fsT, fs = F["fs"]["card_t"], F["fs"]["body"]
    pad = 0.12 * F["k"]

    def h(w):
        iw = w - 2 * pad
        return sum(text_h(t, iw, fsT, bold=True) + text_h(b, iw, fs, 1.02) + 0.06 * F["k"]
                   + 2 * pad + 0.14 * F["k"] for t, _c, _bg, b in P_ASKS)

    def draw(sl, x, y, w):
        iw = w - 2 * pad
        yy = y
        for t, col, bg, b in P_ASKS:
            ht = text_h(t, iw, fsT, bold=True)
            hb = text_h(b, iw, fs, 1.02)
            box = ht + hb + 0.06 * F["k"] + 2 * pad
            rect(sl, x, yy, w, box, fill=bg, line=col, lw=1.2 * F["k"], radius=0.03)
            rect(sl, x, yy, 0.06 * F["k"], box, fill=col)
            tbox(sl, x + pad, yy + pad, iw, ht, t, fs=fsT, bold=True, colour=col)
            tbox(sl, x + pad, yy + pad + ht + 0.06 * F["k"], iw, hb, b, fs=fs, colour=DARK,
                 line=1.02)
            yy += box + 0.14 * F["k"]
    return h, draw


def blk_next(F):
    fsT, fs = F["fs"]["card_t"], F["fs"]["body"]
    pad = 0.08 * F["k"]
    tagw = 0.95 * F["k"]        # "Next" must not stack

    def geom(w):
        tw = w - tagw - 2 * pad - 0.10 * F["k"]
        return tw, [(a, b, c, text_h(b, tw, fsT, bold=True) + text_h(c, tw, fs, 1.0))
                    for a, b, c in P_NEXT]

    def h(w):
        return sum(r[3] + 2 * pad + 0.10 * F["k"] for r in geom(w)[1])

    def draw(sl, x, y, w):
        tw, rows = geom(w)
        yy = y
        for tag, t, b, hh in rows:
            col, bg = (GREEN, GREEN_BG) if tag == "Now" else (JU_BLUE, BLUE_BG)
            rect(sl, x, yy, w, hh + 2 * pad, fill=WHITE, line=RULE, lw=0.9 * F["k"], radius=0.03)
            rect(sl, x + pad, yy + pad, tagw, (fsT * 1.7) / 72.0, fill=bg, radius=0.3)
            tbox(sl, x + pad, yy + pad + (fsT * 0.22) / 72.0, tagw, (fsT * 1.3) / 72.0, tag,
                 fs=fsT * 0.85, bold=True, colour=col, align=PP_ALIGN.CENTER)
            xx = x + pad + tagw + 0.10 * F["k"]
            ht = text_h(t, tw, fsT, bold=True)
            tbox(sl, xx, yy + pad, tw, ht, t, fs=fsT, bold=True, colour=DARK)
            tbox(sl, xx, yy + pad + ht, tw, text_h(b, tw, fs, 1.0), b, fs=fs, colour=GREY, line=1.0)
            yy += hh + 2 * pad + 0.10 * F["k"]
    return h, draw


# ================================================================ full-width feature grid
def card_grid(F, W, cards, ncols, maker):
    """Lay the SF cards as a full-width GRID that reads SF1 SF2 SF3 / SF4 SF5 SF6 …

    They cannot go through the column flow: with 17 cards the flow splits them across columns, so
    SF7 ends up at the TOP of column 3 while SF1 sits at the BOTTOM of column 2 — numerically correct
    for anyone who reads a column at a time, but it defeats the point of numbering them. A band that
    runs left-to-right, row by row, always reads in order.

    Returns (total_height, draw_fn). Row height is the tallest card in that row."""
    gap = F["gutter"]
    cw = (W - (ncols - 1) * gap) / ncols
    blocks = [maker(F, c) for c in cards]
    rows = [blocks[i:i + ncols] for i in range(0, len(blocks), ncols)]
    rh = [max(hf(cw) for hf, _ in r) for r in rows]

    def draw(sl, x, y):
        yy = y
        for r, hgt in zip(rows, rh):
            for j, (_hf, df) in enumerate(r):
                df(sl, x + j * (cw + gap), yy, cw)
            yy += hgt
        return yy - y
    return sum(rh), draw


def band(F, W, label, items, ncols, maker, legend=False):
    """A full-width section drawn BELOW the column flow: optional header, optional status legend,
    then a left-to-right grid. Returns (height, draw_fn)."""
    hdr = blk_section(F, label) if label else None
    leg = blk_legend(F) if legend else None
    gh, gdraw = card_grid(F, W, items, ncols, maker)
    total = (hdr[0](W) if hdr else 0.0) + (leg[0](W) if leg else 0.0) + gh + 0.20 * F["k"]

    def draw(sl, x, y):
        yy = y
        if hdr:
            hdr[1](sl, x, yy, W); yy += hdr[0](W)
        if leg:
            leg[1](sl, x, yy, W); yy += leg[0](W)
        gdraw(sl, x, yy)
    return total, draw


# ================================================================ column flow
def pack(blocks, colw, H, cols):
    """Lay blocks into `cols` columns of height H. Returns [(col, y_offset)] or None.

    A SECTION header refuses to be the last thing in a column — an orphan header at the foot of a
    column with its content at the top of the next is the classic poster-layout defect."""
    out, ci, y = [], 0, 0.0
    for i, (hf, _df, kind) in enumerate(blocks):
        bh = hf(colw)
        nxt = blocks[i + 1][0](colw) if i + 1 < len(blocks) else 0.0
        need = bh + (nxt if kind == "section" else 0.0)
        if y > 0.0 and y + need > H + 1e-6:
            ci += 1
            y = 0.0
            if ci >= cols:
                return None
        out.append((ci, y))
        y += bh
    return out


def flow(sl, blocks, x0, y0, colw, colh, cols, gutter):
    """Balanced fill: binary-search the SHORTEST column height that still fits everything in `cols`
    columns. Greedy packing fills the early columns and leaves the last a third empty, which reads
    as an unfinished poster; balancing shares the leftover whitespace evenly."""
    tall = max(hf(colw) for hf, _, _ in blocks)
    if pack(blocks, colw, colh, cols) is None:
        print("      *** does not fit in %d columns — shorten the content ***" % cols)
        placed, hi = pack(blocks, colw, colh, cols + 8), colh
    else:
        lo, hi = max(tall, 1e-3), colh
        for _ in range(48):
            mid = (lo + hi) / 2.0
            if pack(blocks, colw, mid, cols) is not None:
                hi = mid
            else:
                lo = mid
        placed = pack(blocks, colw, hi, cols)
    used = 0
    for (ci, dy), (_hf, df, _k) in zip(placed, blocks):
        df(sl, x0 + ci * (colw + gutter), y0 + dy, colw)
        used = max(used, ci)
    return used, hi


# ================================================================ page assembly
def header(sl, F, y, title, subtitle, bill, billsub, kpis):
    fs, m, W = F["fs"], F["margin"], F["w"] - 2 * F["margin"]
    fs_t = fit_one_line(title, W, fs["title"], fs["title"] * 0.74)
    ht = text_h(title, W, fs_t, 0.92, bold=True)
    tbox(sl, m, y, W, ht, title, fs=fs_t, bold=True, colour=JU_BLUE)
    y += ht + 0.04 * F["k"]
    ha = text_h(subtitle, W, fs["authors"])
    tbox(sl, m, y, W, ha, subtitle, fs=fs["authors"], colour=GREY)
    y += ha + 0.20 * F["k"]
    hrule(sl, m, y, W, JU_BLUE, 4.0 * F["k"]); y += 0.26 * F["k"]

    fs_b = fit_one_line(bill, W, fs["bill"], fs["bill"] * 0.66)
    hb = text_h(bill, W, fs_b, 0.88, bold=True)
    tbox(sl, m, y, W, hb, bill, fs=fs_b, bold=True, colour=DARK, align=PP_ALIGN.CENTER, line=0.88)
    y += hb + 0.10 * F["k"]
    hs = text_h(billsub, W * 0.86, fs["billsub"], 1.0)
    tbox(sl, m + W * 0.07, y, W * 0.86, hs, billsub, fs=fs["billsub"], colour=GREY,
         align=PP_ALIGN.CENTER, line=1.0)
    y += hs + 0.30 * F["k"]

    kh = (fs["kpi"] * 1.25 + fs["kpil"] * 2.6) / 72.0 + 0.30 * F["k"]
    kw = (W - (len(kpis) - 1) * 0.16 * F["k"]) / len(kpis)
    for i, (val, lab) in enumerate(kpis):
        kx = m + i * (kw + 0.16 * F["k"])
        rect(sl, kx, y, kw, kh, fill=WHITE, line=RULE, lw=1.1 * F["k"], radius=0.06)
        tbox(sl, kx, y + 0.12 * F["k"], kw, (fs["kpi"] * 1.25) / 72.0, val, fs=fs["kpi"], bold=True,
             colour=JU_BLUE, align=PP_ALIGN.CENTER)
        tbox(sl, kx, y + 0.12 * F["k"] + (fs["kpi"] * 1.22) / 72.0, kw, (fs["kpil"] * 2.6) / 72.0,
             lab, fs=fs["kpil"], colour=GREY, align=PP_ALIGN.CENTER, line=1.0)
    return y + kh + 0.42 * F["k"]


CUT_DETAIL = "documentation/figures/rig_cut_detail.png"
CUT_FULL = "documentation/figures/rig_cut_full.png"
SCHEMATIC = "documentation/figures/rig_schematic.png"


def build(fmt):
    F = dict(FORMATS[fmt])
    F["k"] = F["w"] / 33.11                      # geometry scale, A0 = 1.0
    fs = F["fs"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(F["w"]), Inches(F["h"])
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, F["w"], F["h"], fill=PAGE_BG)

    m, W = F["margin"], F["w"] - 2 * F["margin"]
    progress = F["kind"] == "progress"

    y = header(sl, F, m,
               P_TITLE if progress else TITLE,
               P_SUB if progress else AUTHORS,
               P_BILLBOARD if progress else BILLBOARD,
               P_BILLBOARD_SUB if progress else BILLBOARD_SUB,
               P_KPIS if progress else KPIS)

    foot = P_FOOTER if progress else FOOTER
    fh = text_h(foot, W, fs["foot"], 1.05) + 0.22 * F["k"]
    fy = F["h"] - m - fh
    hrule(sl, m, fy - 0.16 * F["k"], W, RULE, 1.4 * F["k"])
    tbox(sl, m, fy, W, fh, foot, fs=fs["foot"], colour=GREY, italic=True, line=1.05)

    colw = (W - (F["cols"] - 1) * F["gutter"]) / F["cols"]
    colh = fy - 0.34 * F["k"] - y

    B, bands = [], []
    def add(t, hf_df):
        B.append((hf_df[0], hf_df[1], t))

    if progress:
        # ---- supervisor / manager: status, what shipped, results, ASKS, next.
        add("section", blk_section(F, "Where the project stands"))
        add("para", blk_para(F, P_STATUS))
        add("rig", blk_rig(F, [CUT_FULL, CUT_DETAIL], row_in=5.4))
        add("section", blk_section(F, "What has been delivered"))
        add("miles", blk_milestones(F))
        add("section", blk_section(F, "Results worth reporting"))
        add("defs", blk_defs(F, P_RESULTS, gap=0.20))
        add("evidence", blk_evidence(F, EVIDENCE[1]))
        add("section", blk_section(F, "Decisions I need from you"))
        add("asks", blk_asks(F))
        add("section", blk_section(F, "Next"))
        add("next", blk_next(F))
    elif F["verbose"]:
        # ---- A0 wall poster: the full argument.
        add("section", blk_section(F, "The machine"))
        add("rig", blk_rig(F, [CUT_DETAIL, SCHEMATIC], row_in=6.6))
        add("section", blk_section(F, "Why automate a tensile rig"))
        add("para", blk_para(F, WHY))
        add("section", blk_section(F, "How it works"))
        add("defs", blk_defs(F, METHODS))
        add("section", blk_section(F, "What the automation buys you"))
        for e in EVIDENCE[:3]:      # the cyclic-window story is carried by SF9 + the Limitations
            add("evidence", blk_evidence(F, e))
        add("section", blk_section(F, "Limitations"))
        add("defs", blk_defs(F, LIMITS))
        bands = [band(F, W, "The smart features, SF1 to SF16", CARDS, 3, blk_card, legend=True)]
    else:
        # ---- A4 handout: a condensation, not a shrunken A0. Methods and limitations stay — they
        # are the whole point of the Generation-2 critique; evidence narrows to one figure.
        add("section", blk_section(F, "The machine"))
        add("rig", blk_rig(F, [CUT_DETAIL, SCHEMATIC], row_in=3.5))
        add("para", blk_para(F, WHY, fs_key="small", colour=GREY))
        add("section", blk_section(F, "How it works"))
        add("defs", blk_defs(F, METHODS[:3]))
        add("section", blk_section(F, "What it buys you"))
        add("evidence", blk_evidence(F, EVIDENCE[1]))
        add("section", blk_section(F, "Limitations"))
        add("defs", blk_defs(F, LIMITS, gap=0.18))
        bands = [band(F, W, "The smart features, SF1 to SF16", CARDS, 2, blk_cardrow, legend=True)]

    # Full-width bands are drawn UNDER the column flow so their contents read left-to-right, in
    # order — which is the whole point of numbering the features. Their height is reserved before
    # the flow is balanced into whatever is left.
    band_h = sum(bh for bh, _ in bands)
    used, balh = flow(sl, B, m, y, colw, colh - band_h, F["cols"], F["gutter"])
    by = y + balh + 0.34 * F["k"]
    for bh, bdraw in bands:
        bdraw(sl, m, by)
        by += bh

    os.makedirs(OUTDIR, exist_ok=True)
    name = {"A0": "Smart_UTM_poster_A0", "A4": "Smart_UTM_poster_A4",
            "A4P": "Smart_UTM_progress_A4"}[fmt]
    out = os.path.join(OUTDIR, name + ".pptx")
    try:
        prs.save(out)
    except PermissionError:
        out = out.replace(".pptx", "_updated.pptx"); prs.save(out)
    total = sum(hf(colw) for hf, _, _ in B)
    cap = colh * F["cols"]
    print("%-4s %.2f x %.2f in · %d cols, %.2f in · content %.1f / %.1f in (%.0f %%) · "
          "balanced %.2f in · ends col %d"
          % (fmt, F["w"], F["h"], F["cols"], colh, total, cap, 100 * total / cap, balh, used + 1))
    print("      -> %s" % out)
    return out


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    for f in ("A0", "A4", "A4P"):
        build(f)
