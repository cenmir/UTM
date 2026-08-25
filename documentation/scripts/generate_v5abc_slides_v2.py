import os as _os  # [doc-folder] run from repo root so plot PNGs & Software/ resolve
_os.chdir(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..', '..'))
"""V5 / V5b / V5c repeatability comparison deck — V2 (50% infill, LED off).
V2 uses the re-run V5c (S2): crosshead position correctly zeroed, but the run was
executed at 0.2 mm/s (vs 0.1 mm/s for V5/V5b) — a strain-rate confound flagged
throughout. Separate output file; the original V1 deck is left untouched.

7 slides, JU template (pages 134-140):
  134 Overview + repeatability headline + rate-confound flag
  135 Stress-strain overlay
  136 Load vs time
  137 Stress vs displacement
  138 Gauge strain vs displacement
  139 Full comparison table (values + %-deviation vs V5)
  140 Offset-factor repeatability + verdict
Output: V5abc_comparison_slides_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

DARK_GREEN = RGBColor(0x1F, 0x3F, 0x2F)
HEADER_GREEN = RGBColor(0x14, 0x3D, 0x2F)
LIGHT_BLUE = RGBColor(0xE7, 0xF1, 0xF8)
GREEN_PASS = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_WARN = RGBColor(0xFF, 0xF3, 0xCD)
ORANGE_FLAG = RGBColor(0xFF, 0xD9, 0xB3)
GREY_TEXT = RGBColor(0x70, 0x70, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
JU_BLUE_BAR = RGBColor(0xA9, 0xD6, 0xEF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, x, y, w, h, text, *, fs=14, bold=False, italic=False, colour=BLACK,
       align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = colour
    return box


def title(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = JU_BLUE_BAR; bar.line.fill.background()
    tb(slide, 0.5, 0.55, 12.5, 0.9, text, fs=30, font="Calibri Light")


def header(slide, x, y, w, text):
    tb(slide, x, y, w, 0.4, text, fs=16, colour=GREY_TEXT)


def cell_bg(cell, rgb):
    cell.fill.solid(); cell.fill.fore_color.rgb = rgb


def cell_txt(cell, text, *, fs=10, bold=False, colour=BLACK, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = colour; r.font.name = "Calibri"


def table(slide, x, y, w, h, data, *, hbg=HEADER_GREEN, hfg=WHITE, cw=None, hf=10, bf=10, ov=None):
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = shp.table
    if cw:
        tot = sum(cw)
        for i, f in enumerate(cw):
            t.columns[i].width = Inches(w*f/tot)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c); text = data[r][c]
            if r == 0:
                cell_bg(cell, hbg); cell_txt(cell, text, fs=hf, bold=True, colour=hfg)
            else:
                cell_txt(cell, text, fs=bf)
            if ov and (r, c) in ov:
                o = ov[(r, c)]
                if 'bg' in o: cell_bg(cell, o['bg'])
                cell_txt(cell, o.get('text', text), fs=bf, bold=o.get('bold', False),
                         colour=o.get('colour', BLACK))
    return shp


def kpi(slide, x, y, w, label, value, *, fill=LIGHT_BLUE, vcol=DARK_GREEN, h=0.98, vfs=18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.text = ""
    r1 = p1.add_run(); r1.text = value
    r1.font.size = Pt(vfs); r1.font.bold = True; r1.font.color.rgb = vcol; r1.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(9.5); r2.font.color.rgb = BLACK; r2.font.name = "Calibri"
    return box


def banner(slide, x, y, w, h, text, *, fill=GREEN_PASS, fg=DARK_GREEN, fs=13):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return box


def pageno(slide, n):
    tb(slide, 0.5, 7.05, 0.7, 0.3, str(n), fs=10, colour=GREY_TEXT)


def ju(slide):
    tb(slide, 12.2, 0.3, 1.0, 0.3, "JÖNKÖPING UNIVERSITY", fs=8, bold=True,
       colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    tb(slide, 0.6, 7.0, 12.1, 0.4, text, fs=12, italic=True, colour=GREY_TEXT)


# ================= SLIDE 1 — Overview + repeatability headline =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5 / V5b / V5c REPEATABILITY (v2) — 50 % INFILL")

header(s, 0.5, 1.3, 6.2, "Three specimens — V5c re-run (position zeroed)")
setup = [
    ["Test", "Spec.", "Preload", "Rate", "Role"],
    ["V5", "S4", "463 N", "0.1 mm/s", "baseline (8.6.20)"],
    ["V5b", "S3", "472 N", "0.1 mm/s", "repeat"],
    ["V5c", "S2", "471 N", "0.2 mm/s", "repeat (re-run, 2x rate)"],
]
ovs = {(3, c): {'bg': ORANGE_FLAG} for c in range(5)}
table(s, 0.5, 1.7, 6.2, 1.5, setup, cw=[0.7, 0.7, 1.0, 1.0, 2.1], hf=11, bf=10.5, ov=ovs)
tb(s, 0.5, 3.28, 6.2, 0.8,
   "All: fresh PLA, 50 % infill (grid), LED off, spray markers, 80 mm² nominal, "
   "preload anchored from post-fracture zero. V5c re-run fixed the displacement "
   "offset (δ zeroed) but was pulled at 0.2 mm/s.", fs=10.5, italic=True, colour=GREY_TEXT)

# KPI tiles (headline repeatability)
kpi(s, 0.5, 4.2, 1.85, "UTS  (CV 2.4 %)", "22.3 MPa", fill=GREEN_PASS)
kpi(s, 2.48, 4.2, 1.85, "peak force (CV 2.4 %)", "1 788 N", fill=GREEN_PASS)
kpi(s, 4.46, 4.2, 1.85, "DIC tracking", "99.7 %", fill=GREEN_PASS)
kpi(s, 0.5, 5.35, 1.85, "σ_y  (CV 9.3 %)", "20.7 MPa", fill=YELLOW_WARN)
kpi(s, 2.48, 5.35, 1.85, "E  (CV 13.6 %)", "1.33 GPa", fill=YELLOW_WARN)
kpi(s, 4.46, 5.35, 1.85, "ε_f  (CV 19.5 %)", "0.023", fill=YELLOW_WARN)

header(s, 6.9, 1.3, 6.0, "Repeatability across the 3 specimens")
rep = [
    ["Metric", "Mean", "Spread", "CV"],
    ["UTS", "22.34 MPa", "4.3 %", "2.4 %"],
    ["Peak force", "1788 N", "4.3 %", "2.4 %"],
    ["σ_y (0.2 % offset)", "20.69 MPa", "18.1 %", "9.3 %"],
    ["Elastic modulus E", "1.33 GPa", "25.2 %", "13.6 %"],
    ["Failure strain ε_f", "0.023", "35.5 %", "19.5 %"],
]
ov = {(1, c): {'bg': GREEN_PASS} for c in range(4)}
ov.update({(2, c): {'bg': GREEN_PASS} for c in range(4)})
ov.update({(4, c): {'bg': YELLOW_WARN} for c in range(4)})
ov.update({(5, c): {'bg': YELLOW_WARN} for c in range(4)})
table(s, 6.9, 1.7, 6.0, 2.3, rep, cw=[1.7, 1.2, 1.0, 1.0], hf=11, bf=11, ov=ov)
tb(s, 6.9, 4.15, 6.0, 1.05,
   "Strength (UTS, peak force) stays tight — CV ≈ 2.4 %. ε_f, toughness and gauge-share "
   "scatter is dominated by V5c's 2× strain rate (PLA stiffens and embrittles at higher "
   "rate → lower ε_f), plus the short-window E slope fit.",
   fs=11, colour=BLACK)
banner(s, 6.9, 5.35, 6.0, 1.0,
       "RATE CONFOUND: V5c ran at 0.2 mm/s (V5/V5b at 0.1). Strength conclusions hold; "
       "ε_f / ductility are NOT rate-matched — re-run V5c at 0.1 mm/s for a clean triplet.",
       fill=ORANGE_FLAG, fg=BLACK, fs=11)
banner(s, 0.5, 6.45, 12.4, 0.5,
       "Strength repeatable (CV 2.4 %); UTS offset ×1.43 holds. ε_f differences are largely the 0.2 mm/s rate, not material scatter.")
pageno(s, 134)

# ================= SLIDE 2 — Stress-strain overlay + noise reason =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — STRESS-STRAIN OVERLAY")
s.shapes.add_picture("images/V5/V5abc_stress_strain.png", Inches(0.35), Inches(1.55), width=Inches(8.5))
header(s, 9.0, 1.4, 4.1, "Reading the overlay")
tb(s, 9.0, 1.85, 4.15, 4.8,
   "• Elastic slope and the ~22–23 MPa plateau overlay closely across all three.\n\n"
   "• V5c (S2, red) reaches UTS at a LOWER strain and fractures earlier "
   "(ε_f 0.0175 vs ~0.025) — consistent with its 2× crosshead rate (PLA is "
   "rate-stiffening / embrittling).\n\n"
   "• Fine ripple = measurement noise, not material: DIC centroid jitter on the "
   "strain axis (floor ≈ 1.4×10⁻⁵) and load-cell ADC noise on the stress axis "
   "(<6 % of the 29.4 kN range).\n\n"
   "• None of it shifts results: E is a slope fit, σ_y the offset construction, UTS a peak.",
   fs=11, colour=BLACK)
footer(s, "Curves coincide on slope and peak (~22–23 MPa). V5c fractures sooner at higher rate (ε_f 0.0175 vs ~0.025) — "
          "a strain-rate effect, not an instrument difference.")
pageno(s, 135)

# ================= SLIDE 3 — Load vs time overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — LOAD vs TIME")
s.shapes.add_picture("images/V5/V5abc_load_time.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "Aligned to ramp start. V5/V5b at 0.1 mm/s, V5c at 0.2 mm/s → V5c ramps ~2× faster and fractures "
          "sooner. Peak force consistent: 1759–1836 N (spread 4.3 %); preload→peak net pull Δ ≈ 1320 N.")
pageno(s, 136)

# ================= SLIDE 4 — Stress vs displacement overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — STRESS vs DISPLACEMENT")
s.shapes.add_picture("images/V5/V5abc_stress_disp.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "Gauge share of crosshead travel: V5 52 %, V5b 53 %, V5c 33 %. V5c's lower share follows its "
          "smaller gauge stretch (ε_f 0.0175 → 1.40 mm) over a larger 4.27 mm travel — partly the higher rate.")
pageno(s, 137)

# ================= SLIDE 5 — Cauchy strain vs displacement overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — GAUGE STRAIN vs DISPLACEMENT")
s.shapes.add_picture("images/V5/V5abc_strain_disp.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "All three sit below the dotted ‘all travel → gauge’ line (rig compliance). V5/V5b reach ε_f ≈ 0.025; "
          "V5c stops at 0.0175 — it fractured earlier at the higher 0.2 mm/s rate.")
pageno(s, 138)

# ================= SLIDE 6 — Full comparison table =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — PARAMETER COMPARISON")
header(s, 0.5, 1.3, 8.0, "Values and % deviation relative to V5 (S4)")
comp = [
    ["Parameter", "V5 (S4)", "V5b (S3)", "V5c (S2)", "Δ% V5b", "Δ% V5c"],
    ["Crosshead rate (mm/s)", "0.10", "0.10", "0.20", "—", "+100"],
    ["Elastic modulus E (GPa)", "1.54", "1.25", "1.20", "−18.4", "−21.9"],
    ["Yield σ_y 0.2 % (MPa)", "19.04", "20.25", "22.79", "+6.4", "+19.7"],
    ["UTS (MPa)", "22.09", "21.99", "22.95", "−0.4", "+3.9"],
    ["Peak true force (N)", "1767", "1759", "1836", "−0.4", "+3.9"],
    ["Failure strain ε_f", "0.0246", "0.0255", "0.0175", "+3.6", "−28.9"],
    ["Toughness (kJ/m³)", "462", "488", "344", "+5.6", "−25.6"],
    ["Gauge share of travel (%)", "51.7", "52.9", "32.8", "+2.3", "−36.6"],
    ["Pull duration (s)", "38.4", "38.7", "21.6", "+0.9", "−43.7"],
    ["DIC tracking (%)", "99.8", "99.7", "99.6", "−0.1", "−0.1"],
]
ovc = {}
for c in range(6):
    ovc[(1, c)] = {'bg': ORANGE_FLAG, 'bold': c == 0}   # rate (confound)
    ovc[(4, c)] = {'bg': GREEN_PASS, 'bold': c == 0}     # UTS
    ovc[(5, c)] = {'bg': GREEN_PASS}                     # peak force
    ovc[(2, c)] = {'bg': YELLOW_WARN}                    # E
    ovc[(6, c)] = {'bg': YELLOW_WARN}                    # eps_f
table(s, 0.5, 1.75, 9.4, 4.55, comp, cw=[2.6, 1.0, 1.0, 1.0, 0.95, 0.95], hf=10.5, bf=10.5, ov=ovc)

header(s, 10.15, 1.3, 2.9, "Reading it")
tb(s, 10.15, 1.75, 3.0, 4.6,
   "• UTS & peak force: within ±4 % — strength reproduces.\n\n"
   "• Orange row = the 0.2 mm/s confound. It drives the V5c deltas below:\n"
   "  – σ_y ↑ (rate stiffening)\n"
   "  – ε_f, toughness, duration, gauge-share ↓ (earlier, more brittle fracture).\n\n"
   "• E scatters (short-window slope fit), independent of rate.\n\n"
   "• DIC tracked 99.6 % even at the higher rate.",
   fs=10.5, colour=BLACK)
footer(s, "Green = repeatable strength. Orange = strain-rate confound (V5c 0.2 mm/s). Yellow = scatter (E fit / rate-driven ε_f).")
pageno(s, 139)

# ================= SLIDE 7 — Offset factor + verdict =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: INFILL OFFSET FACTOR — REPEATABILITY")
s.shapes.add_picture("images/V5/V5abc_offset.png", Inches(0.4), Inches(1.5), width=Inches(7.4))

header(s, 8.1, 1.3, 5.0, "Offset factor  k = Chacón_min / measured")
off = [
    ["Property", "V5", "V5b", "V5c", "mean"],
    ["E", "×1.95", "×2.39", "×2.50", "×2.28"],
    ["σ_y", "×1.58", "×1.48", "×1.32", "×1.46"],
    ["UTS", "×1.45", "×1.45", "×1.39", "×1.43"],
]
ovo = {(3, c): {'bg': GREEN_PASS, 'bold': True} for c in range(5)}
table(s, 8.1, 1.72, 5.0, 1.5, off, cw=[1.1, 0.9, 0.9, 0.9, 0.9], hf=11, bf=11, ov=ovo)

tb(s, 8.1, 3.35, 5.0, 1.0,
   "UTS knock-down is repeatable (×1.43, all 3 within ±0.06); σ_y tight (×1.3–1.6). "
   "E scatters (×1.95–2.50) — it tracks the noisy short-window modulus fit.",
   fs=10.5, colour=BLACK)

banner(s, 8.1, 4.4, 5.0, 1.15,
       "No SINGLE all-property k now: V5c's E demands ≥2.50 but its σ_y caps ≤2.19 "
       "(window empty). Driven by the noisy E fit + V5c's 0.2 mm/s raising σ_y. "
       "The robust result is the strength offset ≈ ×1.4–1.5.",
       fill=YELLOW_WARN, fg=BLACK, fs=11)
tb(s, 8.1, 5.62, 5.0, 0.7,
   "(V1 deck reported a common k ≈ 2.4 from the earlier V5c; the rate-confounded re-run "
   "collapses that window — strength offset is the stable takeaway.)",
   fs=9.5, italic=True, colour=GREY_TEXT)

banner(s, 0.5, 6.5, 12.6, 0.6,
       "Strength offset REPEATABLE (×1.43–1.5). Re-run V5c at 0.1 mm/s for a rate-matched ε_f triplet. "
       "Next: V5d/e (LED on), V6 (100 % infill, expect k ≈ 1).")
pageno(s, 140)

try:
    prs.save("documentation/V5abc_comparison_slides_v2.pptx")
    print("Saved: V5abc_comparison_slides_v2.pptx (7 slides, pages 134-140)")
except PermissionError:
    prs.save("documentation/V5abc_comparison_slides_v2_updated.pptx")
    print("Locked. Saved: V5abc_comparison_slides_v2_updated.pptx")
