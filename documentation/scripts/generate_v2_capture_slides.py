"""V2 / S24 capture-feature validation deck — 2026-08-14, 100 % infill, spray markers.

    python documentation/scripts/v2_capture_plots.py          # plots first
    python documentation/scripts/generate_v2_capture_slides.py

One rig run exercised five features at once. This deck reports what each of them did, and the
material finding that fell out of it: the run is a clean PASS for the capture pipeline and a FAIL
for the fracture point the report picked.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1F, 0x6F, 0xB4)
GREEN = RGBColor(0x2F, 0x9E, 0x44)
DARKGREEN = RGBColor(0x1B, 0x5E, 0x2A)
RED = RGBColor(0xC0, 0x39, 0x2B)
DARKRED = RGBColor(0x7A, 0x24, 0x1B)
AMBER = RGBColor(0xD2, 0x99, 0x22)
DARKAMBER = RGBColor(0x8A, 0x63, 0x10)
PALEGREEN = RGBColor(0xE8, 0xF5, 0xEA)
PALERED = RGBColor(0xFB, 0xEC, 0xEA)
PALEAMBER = RGBColor(0xFD, 0xF6, 0xE3)
PALEGREY = RGBColor(0xF4, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def sl():
    return prs.slides.add_slide(BLANK)


def tb(s, x, y, w, h, text, *, fs=14, bold=False, italic=False, colour=BLACK,
       align=PP_ALIGN.LEFT, wrap=True, font="Calibri", space=0):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align
        p.space_after = Pt(space)
        for run in p.runs:
            run.font.size = Pt(fs); run.font.bold = bold; run.font.italic = italic
            run.font.color.rgb = colour; run.font.name = font
    return box


def title(s, text, sub=None):
    tb(s, 0.55, 0.38, 12.3, 0.7, text, fs=27, bold=True)
    if sub:
        tb(s, 0.55, 1.12, 12.3, 0.45, sub, fs=13.5, colour=GREY)


def rect(s, x, y, w, h, fill, line=None, lw=1.25):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def kpi(s, x, y, w, h, value, label, *, fill, edge, fg):
    rect(s, x, y, w, h, fill, edge)
    tb(s, x, y + 0.13, w, 0.5, value, fs=22, bold=True, colour=fg, align=PP_ALIGN.CENTER)
    tb(s, x, y + 0.66, w, 0.3, label, fs=10.5, colour=GREY, align=PP_ALIGN.CENTER)


def banner(s, x, y, w, h, text, *, fill=PALEGREEN, edge=GREEN, fg=DARKGREEN, fs=13.5, bold=True):
    rect(s, x, y, w, h, fill, edge)
    box = s.shapes.add_textbox(Inches(x + 0.22), Inches(y), Inches(w - 0.44), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(fs); run.font.bold = bold
            run.font.color.rgb = fg; run.font.name = "Calibri"
    return box


def img(s, name, x, y, maxw, maxh):
    p = os.path.join("documentation", "figures", name)
    iw, ih = Image.open(p).size
    sc = min(maxw / iw, maxh / ih)
    w, h = iw * sc, ih * sc
    s.shapes.add_picture(p, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2),
                         Inches(w), Inches(h))


def pageno(s, n):
    tb(s, 12.4, 7.02, 0.6, 0.3, str(n), fs=10.5, colour=GREY, align=PP_ALIGN.RIGHT)


def foot(s, text):
    tb(s, 0.55, 7.02, 10.5, 0.3, text, fs=10.5, colour=GREY)


N = 0


def page(s, note):
    global N
    N += 1
    pageno(s, N); foot(s, note)


RUN = "S24 · V2 · 100 % infill · spray markers · 2026-08-14 12:01 · 188.2 s · 2 135 samples"

# ============================================================== 1 — verdict
s = sl()
title(s, "Frame capture, validated on a real fracture test",
      "One 100 %-infill pull to failure exercised five new features at once — and exposed one "
      "long-standing analysis error.")
kpi(s, 0.55, 1.85, 2.35, 1.05, "1 960", "frames, each sink", fill=PALEGREEN, edge=GREEN, fg=DARKGREEN)
kpi(s, 3.05, 1.85, 2.35, 1.05, "0", "dropped", fill=PALEGREEN, edge=GREEN, fg=DARKGREEN)
kpi(s, 5.55, 1.85, 2.35, 1.05, "99.95 %", "DIC 2/2 markers", fill=PALEGREEN, edge=GREEN, fg=DARKGREEN)
kpi(s, 8.05, 1.85, 2.35, 1.05, "46.5 MPa", "UTS — stands", fill=PALEGREEN, edge=GREEN, fg=DARKGREEN)
kpi(s, 10.55, 1.85, 2.23, 1.05, "17.5 %", "ε_f — WRONG", fill=PALERED, edge=RED, fg=DARKRED)

banner(s, 0.55, 3.15, 12.23, 0.86,
       "CAPTURE PASSES.  Stills, raw video and speckle video each wrote 1 960 frames over 98.3 s "
       "with not one dropped, while the DIC held both markers on 99.95 % of samples all the way "
       "to fracture.")
banner(s, 0.55, 4.15, 12.23, 1.15,
       "THE REPORT'S FRACTURE STRAIN IS NOT A MATERIAL PROPERTY.  ε_f = 17.5 % and toughness "
       "6 774 kJ/m³ come from the broken halves springing apart after the specimen had already "
       "failed. The honest numbers are ε_f = 7.4 % and 2 992 kJ/m³. Strength and stiffness are "
       "unaffected.", fill=PALERED, edge=RED, fg=DARKRED)

rect(s, 0.55, 5.5, 12.23, 1.32, PALEGREY)
tb(s, 0.8, 5.66, 11.7, 1.0,
   "Why this matters more than one specimen:  every 8.6.20 fracture test analysed with the same "
   "detector carries the same error. ε_f and toughness are the two numbers integrated all the way "
   "to the chosen fracture index — so they are the two that move. UTS, σ_y and E are read at or "
   "before the peak and are untouched.", fs=13, colour=BLACK)
page(s, RUN)

# ============================================================== 2 — what ran
s = sl()
title(s, "Five features, one destructive specimen",
      "The run was designed so a single pull would answer every open question from the last two "
      "days of software work.")
rows = [
    ("Frame capture — PNG stills", "PASS",
     "1 960 lossless frames + index.csv, one row per file, 0 dropped", GREEN, PALEGREEN, DARKGREEN),
    ("Frame capture — dual AVI", "PASS",
     "raw + adaptive-speckle written simultaneously, both readable end to end", GREEN, PALEGREEN, DARKGREEN),
    ("GUI unaffected during capture", "PASS",
     "load, DIC and plotting continued at full rate while 2.0 GB was written", GREEN, PALEGREEN, DARKGREEN),
    ("SF11 — CSV ↔ capture link", "PASS",
     "'# Capture:' header written; run.json names the CSV back", GREEN, PALEGREEN, DARKGREEN),
    ("Adaptive speckle under real LEDs", "PASS",
     "markers held to the last frame before fracture — 99.95 % at 2/2", GREEN, PALEGREEN, DARKGREEN),
    ("Fracture-point selection", "FAIL",
     "picked a post-fracture sample; ε_f and toughness both wrong", RED, PALERED, DARKRED),
    ("AVI playback rate", "BUG",
     "header stamped 35 fps, frames arrived at 19.9 — plays 1.76× fast", AMBER, PALEAMBER, DARKAMBER),
]
y = 1.85
for name, verdict, detail, edge, fill, fg in rows:
    rect(s, 0.55, y, 12.23, 0.62, fill, edge, lw=1.0)
    tb(s, 0.78, y + 0.14, 4.3, 0.36, name, fs=13.5, bold=True)
    tb(s, 5.15, y + 0.14, 1.0, 0.36, verdict, fs=13.5, bold=True, colour=fg)
    tb(s, 6.25, y + 0.15, 6.3, 0.36, detail, fs=12, colour=GREY)
    y += 0.72
page(s, RUN)

# ============================================================== 3 — capture timing
s = sl()
title(s, "The writer never fell behind",
      "Every interval between consecutive frames, for the whole 98 s capture. A dropped frame is "
      "not a subtle effect — it doubles an interval, and there is no such interval here.")
img(s, "v2_capture_timing.png", 0.55, 1.75, 12.23, 3.5)
banner(s, 0.55, 5.45, 12.23, 0.82,
       "Median 50.2 ms · p95 53.0 ms · worst 58.2 ms. The slowest frame in the run was 16 % late, "
       "against the 100 % it would take to lose one.")
tb(s, 0.55, 6.45, 12.23, 0.5,
   "This is the measurement that justified the design: cv2 releases the GIL while encoding, so "
   "the writer threads never contend with the GUI for the interpreter.", fs=12, colour=GREY)
page(s, RUN)

# ============================================================== 4 — three sinks
s = sl()
title(s, "Three sinks, one frame count",
      "Stills and both video views were armed together. If any had lagged, its count would differ.")
img(s, "v2_capture_sinks.png", 0.55, 1.7, 8.0, 3.3)
rect(s, 8.9, 1.75, 3.88, 3.2, PALEGREY)
tb(s, 9.12, 1.95, 3.44, 2.9,
   "418 × 2348 px\nMJPG, intra-frame\n\n2.0 GB for a 98 s capture\n— 1.24 GB/min, exactly what "
   "the setup dialog predicted\n\nThe 419→418 crop is the\nodd-width MJPG rule: it crops,\nnever "
   "rescales, so marker\nseparation is preserved", fs=12.5, colour=BLACK, space=3)
banner(s, 0.55, 5.25, 12.23, 1.0,
       "ONE BUG.  Both AVIs carry a 35 fps header — the camera's configured rate — but frames "
       "actually arrived at 19.9 fps. The files play 1.76× too fast, so anything timed off the "
       "video rather than off index.csv will be wrong. The frames themselves are correct.",
       fill=PALEAMBER, edge=AMBER, fg=DARKAMBER)
page(s, RUN)

# ============================================================== 4b — rates unaffected
s = sl()
title(s, "Writing 2.0 GB did not cost the rig a single sample",
      "The requirement when this was specified was that capture must not slow the GUI, the data "
      "gathering rate or the plotting. The run measures it directly — same test, same specimen, "
      "capture on for the middle 98 s.")
img(s, "v2_rate_unaffected.png", 0.55, 1.75, 12.23, 3.4)
banner(s, 0.55, 5.4, 12.23, 0.95,
       "11.39 Hz before · 11.34 Hz during · 11.36 Hz after — a 0.4 % change, smaller than the "
       "variation between the two idle windows. The DIC ran slightly FASTER during capture, "
       "because that is the moving phase where markers are easiest to lock onto.")
tb(s, 0.55, 6.5, 12.23, 0.5,
   "The design reason: OpenCV releases the GIL while encoding, so the PNG and MJPG writer threads "
   "genuinely run beside the Python GUI rather than taking turns with it.", fs=12, colour=GREY)
page(s, RUN)

# ============================================================== 5 — DIC health
s = sl()
title(s, "The DIC held both markers to the end",
      "Adaptive speckle had only ever been proven against synthetic dimming. This is its first "
      "real test under the rig's own LEDs, on a specimen being pulled apart.")
img(s, "v2_dic_health.png", 0.55, 1.8, 12.23, 2.7)
banner(s, 0.55, 4.7, 6.0, 1.5,
       "2 134 of 2 135 samples at 2/2 markers.\n\nThe single exception is one sample at "
       "t = 131.83 s — the fracture instant itself, when one half left the frame.")
rect(s, 6.78, 4.7, 6.0, 1.5, PALEGREY)
tb(s, 7.0, 4.9, 5.56, 1.2,
   "That one dropout is not a defect — it is the detector correctly reporting that the specimen "
   "stopped being one object. It is also the cleanest possible fracture timestamp, and the "
   "analysis should be using it.", fs=12.5, colour=BLACK)
page(s, RUN)

# ============================================================== 6 — the finding
s = sl()
title(s, "The finding: ε_f is measured after the specimen broke",
      "The report places fracture at 17.5 % strain. The last sample at which the specimen was "
      "still in one piece is at 7.4 %.")
img(s, "v2_fracture_point.png", 0.3, 1.6, 8.2, 4.6)
rect(s, 8.7, 1.75, 4.08, 2.05, PALEGREY)
tb(s, 8.92, 1.95, 3.64, 1.75,
   "Between the green dot and the red X the load has already collapsed from 2 986 N to below "
   "zero. Nothing in that span is the material deforming — it is two separated halves and a "
   "relaxing frame.", fs=12.5, colour=BLACK)
banner(s, 8.7, 4.0, 4.08, 1.05,
       "ε_f  17.5 %  →  7.4 %", fill=PALERED, edge=RED, fg=DARKRED, fs=17)
banner(s, 8.7, 5.2, 4.08, 1.05,
       "toughness\n6 774  →  2 992 kJ/m³", fill=PALERED, edge=RED, fg=DARKRED, fs=14)
page(s, RUN)

# ============================================================== 7 — smoking gun
s = sl()
title(s, "What the markers actually did",
      "Marker separation and load, either side of fracture. They break at the same instant, which "
      "is the whole argument.")
img(s, "v2_lpx_jump.png", 0.55, 1.7, 12.23, 3.7)
banner(s, 0.55, 5.6, 12.23, 1.0,
       "The gauge grew 196 px — 10.9 % — in 0.45 s. Through the preceding 90 s of pulling it grew "
       "at 1.26 × 10⁻³ /s. The jump implies 2.63 × 10⁻¹ /s: 209 times faster, at the exact moment "
       "the load fell to zero. PLA does not do that. A snapping specimen does.",
       fill=PALERED, edge=RED, fg=DARKRED)
page(s, RUN)

# ============================================================== 8 — impossibility
s = sl()
title(s, "A check that needs no DIC at all",
      "Strain is a ratio. Multiply it back into millimetres and compare against how far the "
      "machine actually moved.")
img(s, "v2_gauge_impossible.png", 1.3, 1.75, 10.7, 3.9)
banner(s, 0.55, 5.85, 12.23, 0.95,
       "ε_f = 17.5 % on an 80 mm gauge means the gauge alone stretched 14.04 mm. The crosshead "
       "moved 9.35 mm in the entire test — grips, load train, frame and specimen combined. The "
       "gauge cannot out-stretch the machine, so 17.5 % is impossible on arithmetic alone.",
       fill=PALERED, edge=RED, fg=DARKRED)
page(s, RUN)

# ============================================================== 9 — corrected KPIs
s = sl()
title(s, "What changes, and what does not",
      "Only the two quantities integrated all the way to the fracture index move.")
img(s, "v2_kpi_correction.png", 0.9, 1.7, 11.5, 3.4)
rect(s, 0.55, 5.25, 6.0, 1.55, PALEGREEN, GREEN)
tb(s, 0.78, 5.44, 5.54, 1.25,
   "UNCHANGED — read at or before the peak\n\nUTS 46.5 MPa · σ_y 40.5 MPa · E 2.93 GPa\n\n"
   "Sits inside the V6 quintet (46.2 ± 1.2 MPa) and above its E of 2.60 GPa.",
   fs=12.5, colour=DARKGREEN, space=2)
rect(s, 6.78, 5.25, 6.0, 1.55, PALERED, RED)
tb(s, 7.0, 5.44, 5.54, 1.25,
   "CORRECTED — integrated to fracture\n\nε_f  17.5 % → 7.44 %\ntoughness  6 774 → 2 992 kJ/m³ "
   "(2.26× overstated)\n\n7.4 % also sits much closer to the 3–7 % of the V6 quintet.",
   fs=12.5, colour=DARKRED, space=2)
page(s, RUN)

# ============================================================== 10 — actions
s = sl()
title(s, "What to do about it",
      "Three fixes, in the order they matter.")
items = [
    ("1", "Take the fracture index from the load collapse, not the strain",
     "The detector already finds the collapse — the load fell 1 957 N in one 80 ms sample. "
     "ε_f should be read at the last sample BEFORE it, not after. The DIC_Blobs column marks the "
     "same instant independently.", RED, PALERED, DARKRED),
    ("2", "Re-run every 8.6.20 fracture test through the corrected detector",
     "ε_f and toughness for the whole V5/V6 series were computed the same way. The V6 quintet's "
     "3–7 % range suggests those were mostly fine — but 'mostly' is not a number you publish.",
     AMBER, PALEAMBER, DARKAMBER),
    ("3", "Stamp the achieved frame rate into the AVI, not the configured one",
     "One line in the video writer. Until then, time anything off frames/index.csv, which carries "
     "a real PC timestamp per frame and is correct.", AMBER, PALEAMBER, DARKAMBER),
]
y = 1.9
for num, head, body, edge, fill, fg in items:
    rect(s, 0.55, y, 12.23, 1.35, fill, edge)
    tb(s, 0.8, y + 0.3, 0.5, 0.7, num, fs=26, bold=True, colour=fg)
    tb(s, 1.45, y + 0.16, 11.1, 0.4, head, fs=15, bold=True, colour=fg)
    tb(s, 1.45, y + 0.6, 11.1, 0.7, body, fs=12.5, colour=BLACK)
    y += 1.5
banner(s, 0.55, 6.45, 12.23, 0.62,
       "The capture feature itself needs no changes. It did exactly what it was built to do, on "
       "the first destructive test it was asked to survive.")
page(s, RUN)

OUT = os.path.join("documentation", "decks", "V2_capture_validation.pptx")
prs.save(OUT)
print(f"wrote {OUT}  ({N} slides)")
