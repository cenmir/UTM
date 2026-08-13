import os as _os  # [doc-folder] run from repo root so the test CSVs and PNGs resolve
_os.chdir(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..'))
"""How Young's modulus (E) is calculated — a 2-slide explainer in plain language.

Slide 1  the idea: E is the steepness of the straight bit — rise divided by run, drawn on the
         real measurements from V6d (S11, 100 % infill).
Slide 2  the real formula the code uses (least squares over all 101 points), worked through with
         the actual numbers from that test, and where on the curve it is measured.

Every number is RECOMPUTED from the CSV by e_explainer_plots, never typed in, so the slides cannot
drift from what utm_analysis.analyze() actually does.

Output: documentation/E_modulus_explained.pptx
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

sys.path.insert(0, "documentation")
import e_explainer_plots as EP          # recomputes E, c1, R2, the window and the figures

# ---- house palette, same as generate_v6a_slides.py -----------------------------------------
DARK_GREEN = RGBColor(0x1F, 0x3F, 0x2F); HEADER_GREEN = RGBColor(0x14, 0x3D, 0x2F)
LIGHT_BLUE = RGBColor(0xE7, 0xF1, 0xF8); GREEN_PASS = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_WARN = RGBColor(0xFF, 0xF3, 0xCD)
GREY_TEXT = RGBColor(0x70, 0x70, 0x70); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00); JU_BLUE_BAR = RGBColor(0xA9, 0xD6, 0xEF)
FLOW_BLUE = RGBColor(0x1F, 0x77, 0xB4); FLOW_RED = RGBColor(0xD6, 0x27, 0x28)
GREEN_RUN = RGBColor(0x2A, 0x9D, 0x5C)
CODE_BG = RGBColor(0xF4, 0xF6, 0xF8); CODE_BORDER = RGBColor(0xB8, 0xC2, 0xCC)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, x, y, w, h, text, *, fs=14, bold=False, italic=False, colour=BLACK,
       align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = colour
    return box


def title(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = JU_BLUE_BAR; bar.line.fill.background()
    tb(slide, 0.5, 0.55, 12.5, 0.9, text, fs=30, font="Calibri Light")


def header(slide, x, y, w, text):
    tb(slide, x, y, w, 0.4, text, fs=16, colour=GREY_TEXT)


def banner(slide, x, y, w, h, text, *, fill=GREEN_PASS, fg=DARK_GREEN, fs=13):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return box


def pageno(slide, n):
    tb(slide, 12.45, 7.05, 0.6, 0.3, str(n), fs=10, colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def ju(slide):
    tb(slide, 12.2, 0.3, 1.0, 0.3, "JÖNKÖPING UNIVERSITY", fs=8, bold=True,
       colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    tb(slide, 0.6, 7.0, 11.75, 0.4, text, fs=12, italic=True, colour=GREY_TEXT)


def img_fit(slide, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih)
    w, h = iw * sc, ih * sc
    slide.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y), width=Inches(w))
    return h


def codebox(slide, x, y, w, h, lines, *, fs=12.5):
    """Monospaced panel — the arithmetic has to line up column-wise to be readable."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CODE_BORDER; box.line.width = Pt(1.1)
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.09); tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (txt, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = txt
        r.font.size = Pt(fs); r.font.bold = bold; r.font.name = "Consolas"
        r.font.color.rgb = col
    return box


def bullet(slide, x, y, w, lead, rest, *, fs=13.5, lead_col=FLOW_BLUE):
    """A short bold lead-in followed by plain text, as one wrapped paragraph."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.75))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = lead
    r1.font.size = Pt(fs); r1.font.bold = True; r1.font.name = "Calibri"; r1.font.color.rgb = lead_col
    r2 = p.add_run(); r2.text = rest
    r2.font.size = Pt(fs); r2.font.name = "Calibri"; r2.font.color.rgb = BLACK
    return box


# ============================================================================================
E, C1, R2, N = EP.E, EP.C1, EP.R2, len(EP.WIN)
EPS, SIG = EP.EPS, EP.SIG
X1, X2 = 0.0005, 0.004
Y1, Y2 = E * X1 + C1, E * X2 + C1
RISE, RUN = Y2 - Y1, X2 - X1
ANCHOR, AREA = EP.R["anchor"], EP.AREA
n = N
Se, Ss = sum(EPS), sum(SIG)
See, Ses = sum(e * e for e in EPS), sum(e * s for e, s in zip(EPS, SIG))
NUM, DEN = n * Ses - Se * Ss, n * See - Se * Se

# ===== SLIDE 1 — the idea ====================================================================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "Young's modulus  E  —  how stiff is the material?")
header(s, 0.5, 1.32, 12.4,
       "Pull the specimen. It stretches. E is how steeply the pull has to rise to keep it stretching.")

img_fit(s, "documentation/e_fig_riserun.png", 0.35, 1.80, 7.35, 4.55)

XR, WR = 7.95, 4.95
bullet(s, XR, 1.82, WR, "STRESS  ",
       "= how hard you are pulling ÷ how thick the specimen is.  Units MPa.")
bullet(s, XR, 2.52, WR, "STRAIN  ",
       "= how much the gauge stretched ÷ how long it was.  Just a fraction — no units.")
bullet(s, XR, 3.32, WR, "E  ",
       "= how much the STRESS went up, divided by how much it STRETCHED to get there.")

codebox(s, XR, 4.18, WR, 1.30, [
    (f"RISE  =  {Y2:6.2f}  -  {Y1:5.2f}   =  {RISE:.2f} MPa", True, GREEN_RUN),
    (f"RUN   =  {X2:.4f}  - {X1:.4f}  =  {RUN:.4f}", True, GREEN_RUN),
    ("", False, BLACK),
    (f"E  =  {RISE:.2f} / {RUN:.4f}  =  {E:.0f} MPa", True, FLOW_RED),
], fs=12)

tb(s, XR, 5.62, WR, 0.75,
   f"{E:.0f} MPa = {E/1000:.2f} GPa.  A stiffer material gives a steeper line and a bigger E; "
   "a floppier one gives a shallower line.",
   fs=12.5, colour=GREY_TEXT)

banner(s, 0.5, 6.55, 12.35, 0.42,
       "Steepness of the straight bit = stiffness.  Nothing else on the curve is used for E.",
       fill=LIGHT_BLUE, fg=DARK_GREEN, fs=13)
footer(s, "V6d · S11 · 100 % infill — the real measurements, not an illustration")
pageno(s, 1)

# ===== SLIDE 2 — the real formula ============================================================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "The formula the code actually uses")
header(s, 0.5, 1.32, 12.4,
       f"Two points would do — but {N} were measured, so all {N} are used. That is all "
       "\"least squares\" means.")

img_fit(s, "documentation/e_fig_window.png", 0.35, 1.80, 6.15, 4.60)

XR, WR = 6.85, 6.05
tb(s, XR, 1.78, WR, 0.35, "Best-fit slope through all the points:", fs=13.5, bold=True)
codebox(s, XR, 2.16, WR, 1.02, [
    ("        n·Σ(ε·σ)  −  Σε · Σσ", True, FLOW_RED),
    ("  E =  ─────────────────────────", True, FLOW_RED),
    ("        n·Σ(ε·ε)  −  (Σε)²", True, FLOW_RED),
], fs=13)

tb(s, XR, 3.30, WR, 0.35, f"Add up four columns over the {N} points:", fs=13.5, bold=True)
codebox(s, XR, 3.68, WR, 1.30, [
    (f"n        = {n}", False, BLACK),
    (f"Σε       = {Se:.6f}", False, BLACK),
    (f"Σσ       = {Ss:.4f}  MPa", False, BLACK),
    (f"Σ(ε·ε)   = {See:.9f}", False, BLACK),
    (f"Σ(ε·σ)   = {Ses:.6f}  MPa", False, BLACK),
], fs=12)

tb(s, XR, 5.08, WR, 0.35, "Put the numbers in:", fs=13.5, bold=True)
codebox(s, XR, 5.46, WR, 1.02, [
    (f"top    = {n}·{Ses:.4f} − {Se:.4f}·{Ss:.2f} = {NUM:.4f}", False, BLACK),
    (f"bottom = {n}·{See:.7f} − {Se:.4f}²  = {DEN:.7f}", False, BLACK),
    (f"E = {NUM:.4f} / {DEN:.7f} = {E:.0f} MPa = {E/1000:.2f} GPa", True, FLOW_RED),
], fs=11.5)

banner(s, 0.5, 6.55, 12.35, 0.42,
       f"σ = {E:.0f} × ε + {C1:.2f}      R² = {R2:.4f}      "
       f"the {C1:.2f} MPa offset is the preload the curve starts from, not an error",
       fill=YELLOW_WARN, fg=DARK_GREEN, fs=12.5)
footer(s, f"utm_analysis.analyze() · fit window 0.05 %–0.40 % strain · "
          f"σ = (F + {ANCHOR:.0f} N anchor) / {AREA:.0f} mm²")
pageno(s, 2)

# ===== SLIDE 3 — which line is which =========================================================
def step(slide, x, y, w, h, num, text, *, col):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = col; box.line.width = Pt(1.6)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = f"{num}   "
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.name = "Calibri"; r1.font.color.rgb = col
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(11); r2.font.name = "Calibri"; r2.font.color.rgb = BLACK
    return box


s = prs.slides.add_slide(BLANK); ju(s)
title(s, "The two grey lines on the report's stress–strain graph")
header(s, 0.5, 1.32, 12.4,
       "Both have the SAME slope — the one E you measured. Neither of them is used to compute E.")

img_fit(s, "documentation/e_fig_two_lines.png", 0.85, 1.72, 11.6, 4.40)

SW, GAP = 2.94, 0.18
for i, (num, txt, col) in enumerate([
        ("1", "Fit E first — least squares over the points from 0.05 % to 0.40 % strain.",
         FLOW_RED),
        ("2", "Draw ① with that slope: the short dashed line that hugs the start of the curve.",
         RGBColor(0x33, 0x33, 0x33)),
        ("3", "Copy it and slide it 0.2 % to the right — that copy is ②, the long dotted line.",
         RGBColor(0xB0, 0x30, 0x60)),
        ("4", "Where ② meets the curve is σ_y, the 0.2 % offset yield stress.",
         RGBColor(0x2A, 0x9D, 0x5C))]):
    step(s, 0.55 + i * (SW + GAP), 6.22, SW, 0.62, num, txt, col=col)

footer(s, "The line is not an extrapolation FROM the yield point — the yield point is defined BY "
          "the line. S16 · E 1.88 GPa · σ_y 47.0 MPa")
pageno(s, 3)

OUT = "documentation/E_modulus_explained.pptx"
prs.save(OUT)
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"  E={E:.1f} MPa  c1={C1:.3f}  R2={R2:.5f}  n={N}")
