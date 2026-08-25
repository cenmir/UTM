import os as _os  # [doc-folder] run from repo root so plot PNGs & Software/ resolve
_os.chdir(_os.path.join(_os.path.dirname(_os.path.abspath(__file__)), '..'))
"""V5 / V5b / V5c repeatability comparison deck (50% infill, LED off).
4 slides, JU template (pages 134-137):
  134 Overview + repeatability headline
  135 Stress-strain overlay
  136 Full comparison table (values + %-deviation vs V5)
  137 Offset-factor repeatability + verdict
Output: V5abc_comparison_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

DARK_GREEN = RGBColor(0x1F, 0x3F, 0x2F)
HEADER_GREEN = RGBColor(0x14, 0x3D, 0x2F)
LIGHT_BLUE = RGBColor(0xE7, 0xF1, 0xF8)
GREEN_PASS = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_WARN = RGBColor(0xFF, 0xF3, 0xCD)
GREY_TEXT = RGBColor(0x70, 0x70, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
JU_BLUE_BAR = RGBColor(0xA9, 0xD6, 0xEF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, x, y, w, h, text, *, fs=14, bold=False, italic=False, colour=BLACK,
       align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02); tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = colour
    return box


def title(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.45), Inches(0.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = JU_BLUE_BAR; bar.line.fill.background()
    tb(slide, 0.5, 0.55, 12.5, 0.9, text, fs=30, font="Calibri Light")


def header(slide, x, y, w, text):
    tb(slide, x, y, w, 0.4, text, fs=16, colour=GREY_TEXT)


def cell_bg(cell, rgb):
    cell.fill.solid(); cell.fill.fore_color.rgb = rgb


def cell_txt(cell, text, *, fs=10, bold=False, colour=BLACK, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = colour; r.font.name = "Calibri"


def table(slide, x, y, w, h, data, *, hbg=HEADER_GREEN, hfg=WHITE, cw=None, hf=10, bf=10, ov=None):
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = shp.table
    if cw:
        tot = sum(cw)
        for i, f in enumerate(cw):
            t.columns[i].width = Inches(w*f/tot)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c); text = data[r][c]
            if r == 0:
                cell_bg(cell, hbg); cell_txt(cell, text, fs=hf, bold=True, colour=hfg)
            else:
                cell_txt(cell, text, fs=bf)
            if ov and (r, c) in ov:
                o = ov[(r, c)]
                if 'bg' in o: cell_bg(cell, o['bg'])
                cell_txt(cell, o.get('text', text), fs=bf, bold=o.get('bold', False),
                         colour=o.get('colour', BLACK))
    return shp


def kpi(slide, x, y, w, label, value, *, fill=LIGHT_BLUE, vcol=DARK_GREEN, h=0.98, vfs=18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.text = ""
    r1 = p1.add_run(); r1.text = value
    r1.font.size = Pt(vfs); r1.font.bold = True; r1.font.color.rgb = vcol; r1.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(9.5); r2.font.color.rgb = BLACK; r2.font.name = "Calibri"
    return box


def banner(slide, x, y, w, h, text, *, fill=GREEN_PASS, fg=DARK_GREEN, fs=13):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return box


def pageno(slide, n):
    tb(slide, 0.5, 7.05, 0.7, 0.3, str(n), fs=10, colour=GREY_TEXT)


def ju(slide):
    tb(slide, 12.2, 0.3, 1.0, 0.3, "JÖNKÖPING UNIVERSITY", fs=8, bold=True,
       colour=GREY_TEXT, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    tb(slide, 0.6, 7.0, 12.1, 0.4, text, fs=12, italic=True, colour=GREY_TEXT)


# ================= SLIDE 1 — Overview + repeatability headline =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5 / V5b / V5c REPEATABILITY — 50 % INFILL")

header(s, 0.5, 1.35, 6.0, "Three specimens, identical conditions")
setup = [
    ["Test", "Specimen", "Preload", "Role"],
    ["V5", "S4", "463 N", "baseline (8.6.20)"],
    ["V5b", "S3", "472 N", "repeat"],
    ["V5c", "S2", "466 N", "repeat"],
]
table(s, 0.5, 1.75, 6.0, 1.5, setup, cw=[0.8, 1.0, 1.0, 1.8], hf=11, bf=11)
tb(s, 0.5, 3.35, 6.0, 0.7,
   "All: fresh PLA, 50 % infill (grid pattern), LED off, spray markers, "
   "0.1 mm/s continuous to fracture, 80 mm² nominal area, preload anchored from "
   "post-fracture zero.", fs=11, italic=True, colour=GREY_TEXT)

# KPI tiles (headline repeatability)
kpi(s, 0.5, 4.25, 1.85, "UTS  (CV 0.2 %)", "22.0 MPa", fill=GREEN_PASS)
kpi(s, 2.48, 4.25, 1.85, "peak force (CV 0.2 %)", "1 763 N", fill=GREEN_PASS)
kpi(s, 4.46, 4.25, 1.85, "DIC tracking", "99.7 %", fill=GREEN_PASS)
kpi(s, 0.5, 5.4, 1.85, "σ_y  (CV 3.6 %)", "19.9 MPa")
kpi(s, 2.48, 5.4, 1.85, "E  (CV 10 %)", "1.38 GPa", fill=YELLOW_WARN)
kpi(s, 4.46, 5.4, 1.85, "ε_f  (CV 13 %)", "0.027", fill=YELLOW_WARN)

header(s, 6.8, 1.35, 6.0, "Repeatability across the 3 specimens")
rep = [
    ["Metric", "Mean", "Spread", "CV"],
    ["UTS", "22.03 MPa", "0.4 %", "0.2 %"],
    ["Peak force", "1763 N", "0.4 %", "0.2 %"],
    ["σ_y (0.2 % offset)", "19.86 MPa", "6.4 %", "3.6 %"],
    ["Elastic modulus E", "1.38 GPa", "20.5 %", "10.3 %"],
    ["Failure strain ε_f", "0.027", "23.2 %", "12.6 %"],
]
ov = {(1, c): {'bg': GREEN_PASS} for c in range(4)}
ov.update({(2, c): {'bg': GREEN_PASS} for c in range(4)})
ov.update({(4, c): {'bg': YELLOW_WARN} for c in range(4)})
ov.update({(5, c): {'bg': YELLOW_WARN} for c in range(4)})
table(s, 6.8, 1.75, 6.0, 2.3, rep, cw=[1.7, 1.2, 1.0, 1.0], hf=11, bf=11, ov=ov)
tb(s, 6.8, 4.25, 6.0, 1.0,
   "Strength (UTS, peak force) is outstandingly repeatable — CV < 0.5 %. Modulus E "
   "and ε_f scatter more: E is a slope fit over a short low-strain window (most "
   "noise-sensitive), and S2 (V5c) stretched ~25 % further before fracture.",
   fs=11.5, colour=BLACK)
banner(s, 0.5, 6.35, 12.3, 0.55,
       "Strength repeatable to < 0.5 % across 3 specimens — the 50 % infill response is stable and reproducible.")
pageno(s, 134)

# ================= SLIDE 2 — Stress-strain overlay + noise reason =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — STRESS-STRAIN OVERLAY")
s.shapes.add_picture("images/V5/V5abc_stress_strain.png", Inches(0.35), Inches(1.55), width=Inches(8.5))
header(s, 9.0, 1.4, 4.1, "Why the curves look “noisy”")
tb(s, 9.0, 1.85, 4.15, 4.7,
   "The fine ripple/steps are MEASUREMENT noise, not material:\n\n"
   "• Strain axis (dominant): DIC marker-centroid jitter — the blob centroid wobbles "
   "~0.1–1 px between frames over a ~1675 px gauge → ε noise floor ≈ 1.4×10⁻⁵ "
   "(the small staircase steps).\n\n"
   "• Stress axis: load-cell ADC noise. The 3 t (29.4 kN) cell runs at < 6 % of range, "
   "so ±0.2–1.6 N baseline = only ±0.003–0.02 MPa ripple.\n\n"
   "• Force–DIC timing jitter: per-sample lag varies 0–140 ms (Lag_ms column), smearing "
   "points where the curve is steep.\n\n"
   "None of it shifts the results: E is a slope fit, σ_y the offset construction, UTS a "
   "peak — all average or ignore the scatter.",
   fs=11, colour=BLACK)
footer(s, "Near-coincident curves: identical elastic slope and UTS plateau (~22 MPa); only V5c (S2) "
          "stretches to a higher failure strain (0.031 vs ~0.025).")
pageno(s, 135)

# ================= SLIDE 3 — Load vs time overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — LOAD vs TIME")
s.shapes.add_picture("images/V5/V5abc_load_time.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "Aligned to ramp start (all 0.1 mm/s). Peak force is repeatable to 8 N (1759–1767 N, 0.5 %); "
          "V5c (S2) sustains the peak ~10 s longer before fracture — its extra ductility.")
pageno(s, 136)

# ================= SLIDE 4 — Stress vs displacement overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — STRESS vs DISPLACEMENT")
s.shapes.add_picture("images/V5/V5abc_stress_disp.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "Gauge share of crosshead travel: V5 52 %, V5b 53 %, V5c 50 % (range 2.6 pts) — consistently "
          "~half the travel stretches the gauge, half is rig/grip take-up. V5c needed 4.9 mm travel (vs 3.8) to fracture.")
pageno(s, 137)

# ================= SLIDE 5 — Cauchy strain vs displacement overlay =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — GAUGE STRAIN vs DISPLACEMENT")
s.shapes.add_picture("images/V5/V5abc_strain_disp.png", Inches(1.67), Inches(1.55), width=Inches(10.0))
footer(s, "All three sit well below the dotted ‘all travel → gauge’ line, confirming ~half the travel is rig "
          "compliance. Slopes overlay closely; V5c continues to 0.031 strain at 4.9 mm before fracture.")
pageno(s, 138)

# ================= SLIDE 3 — Full comparison table =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: V5/V5b/V5c — PARAMETER COMPARISON")
header(s, 0.5, 1.35, 8.0, "Values and % deviation relative to V5 (S4)")
comp = [
    ["Parameter", "V5 (S4)", "V5b (S3)", "V5c (S2)", "Δ% V5b", "Δ% V5c"],
    ["Elastic modulus E (GPa)", "1.54", "1.25", "1.36", "−18.4", "−11.5"],
    ["E fit R²", "0.994", "0.995", "0.937", "—", "—"],
    ["Yield σ_y 0.2 % (MPa)", "19.04", "20.25", "20.30", "+6.4", "+6.6"],
    ["UTS (MPa)", "22.09", "21.99", "22.02", "−0.4", "−0.3"],
    ["Peak true force (N)", "1767", "1759", "1762", "−0.4", "−0.3"],
    ["Failure strain ε_f", "0.0246", "0.0255", "0.0309", "+3.6", "+25.5"],
    ["Fracture stress (MPa)", "21.75", "21.22", "20.62", "−2.4", "−5.2"],
    ["Toughness (kJ/m³)", "462", "488", "692", "+5.6", "+49.7"],
    ["DIC tracking (%)", "99.8", "99.7", "99.7", "−0.1", "−0.1"],
]
ovc = {}
for c in range(6):
    ovc[(4, c)] = {'bg': GREEN_PASS, 'bold': c == 0}   # UTS
    ovc[(5, c)] = {'bg': GREEN_PASS}                    # peak force
    ovc[(1, c)] = {'bg': YELLOW_WARN}                   # E
    ovc[(6, c)] = {'bg': YELLOW_WARN}                   # eps_f
table(s, 0.5, 1.8, 9.4, 4.3, comp, cw=[2.5, 1.0, 1.0, 1.0, 0.95, 0.95], hf=10.5, bf=11, ov=ovc)

header(s, 10.15, 1.35, 2.9, "Reading it")
tb(s, 10.15, 1.8, 3.0, 4.2,
   "• UTS & peak force: ±0.4 % — essentially identical.\n\n"
   "• σ_y: +6 % (both repeats higher), tight band.\n\n"
   "• E scatters −12…−18 %: short-window slope fit; V5c's elastic R² = 0.94 "
   "(more DIC scatter) drives its E down.\n\n"
   "• V5c (S2) failed at higher ε_f → larger toughness & post-UTS softening. "
   "Specimen ductility variation, not a rig effect (DIC tracked 99.7 %).",
   fs=11, colour=BLACK)
footer(s, "Green = excellent repeatability (strength). Yellow = scatter sources (modulus fit, S2 ductility).")
pageno(s, 139)

# ================= SLIDE 4 — Offset factor + verdict =================
s = prs.slides.add_slide(BLANK); ju(s)
title(s, "8.6.20: INFILL OFFSET FACTOR — REPEATABILITY")
s.shapes.add_picture("images/V5/V5abc_offset.png", Inches(0.4), Inches(1.5), width=Inches(7.4))

header(s, 8.1, 1.35, 5.0, "Offset factor  k = Chacón_min / measured")
off = [
    ["Property", "V5", "V5b", "V5c", "mean"],
    ["E", "×1.95", "×2.39", "×2.21", "×2.18"],
    ["σ_y", "×1.58", "×1.48", "×1.48", "×1.51"],
    ["UTS", "×1.45", "×1.45", "×1.45", "×1.45"],
]
ovo = {(3, c): {'bg': GREEN_PASS, 'bold': True} for c in range(5)}
table(s, 8.1, 1.78, 5.0, 1.5, off, cw=[1.1, 0.9, 0.9, 0.9, 0.9], hf=11, bf=11, ov=ovo)

tb(s, 8.1, 3.45, 5.0, 0.9,
   "UTS knock-down is identical (×1.45) on all 3 specimens; σ_y tight (×1.5). "
   "E scatters (×1.95–2.39) — it tracks the noisy modulus fit.",
   fs=11, colour=BLACK)

banner(s, 8.1, 4.45, 5.0, 1.0,
       "Single common factor for all 3 specimens: k ≈ 2.4 (window [2.39, 2.46]). "
       "Apply k ≈ 2.4 → E, σ_y, UTS all land in Chacón ranges.",
       fill=LIGHT_BLUE, fg=DARK_GREEN, fs=12)
tb(s, 8.1, 5.55, 5.0, 0.8,
   "(V5 alone gave [1.95, 2.63], k = 2.0; adding V5b/V5c tightened the window — "
   "now bound from below by modulus.)", fs=10, italic=True, colour=GREY_TEXT)

banner(s, 0.5, 6.5, 12.6, 0.6,
       "The infill offset is REPEATABLE — a stable, calibratable knock-down (strength ×1.45–1.5). "
       "Next: V5d/e (LED on) and V6 (100 % infill, expect k ≈ 1).")
pageno(s, 140)

try:
    prs.save("documentation/V5abc_comparison_slides.pptx")
    print("Saved: V5abc_comparison_slides.pptx (7 slides, pages 134-140)")
except PermissionError:
    prs.save("documentation/V5abc_comparison_slides_updated.pptx")
    print("Original locked (open in PowerPoint). Saved: V5abc_comparison_slides_updated.pptx")
